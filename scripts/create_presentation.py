from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color palette
BG_DARK = RGBColor(0x0F, 0x17, 0x2A)
BG_CARD = RGBColor(0x1A, 0x25, 0x3C)
ACCENT = RGBColor(0x38, 0xBD, 0xF8)
ACCENT2 = RGBColor(0x22, 0xD3, 0xEE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x94, 0xA3, 0xB8)
LIGHT = RGBColor(0xCB, 0xD5, 0xE1)
ORANGE = RGBColor(0xFB, 0x92, 0x3C)
GREEN = RGBColor(0x4A, 0xDE, 0x80)
RED = RGBColor(0xF8, 0x71, 0x71)
YELLOW = RGBColor(0xFA, 0xCC, 0x15)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    # Round corners
    shape.adjustments[0] = 0.05
    return shape


def add_text(slide, left, top, width, height, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name="맑은 고딕"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox


def add_bullet_list(slide, left, top, width, height, items, size=16, color=LIGHT, bullet_color=ACCENT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(8)
        # Bullet
        run1 = p.add_run()
        run1.text = "▸ "
        run1.font.size = Pt(size)
        run1.font.color.rgb = bullet_color
        run1.font.name = "맑은 고딕"
        # Text
        run2 = p.add_run()
        run2.text = item
        run2.font.size = Pt(size)
        run2.font.color.rgb = color
        run2.font.name = "맑은 고딕"
    return txBox


# ============================================================
# SLIDE 1: Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, BG_DARK)

# Accent line
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(3.1), Inches(13.333), Pt(3))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()

add_text(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1),
         "GLANCE", size=54, color=ACCENT, bold=True)
add_text(slide, Inches(1.5), Inches(2.45), Inches(10), Inches(0.7),
         "AI Code Generation 기반 현장 밸브 관리 프로덕트", size=28, color=WHITE, bold=True)

add_text(slide, Inches(1.5), Inches(3.6), Inches(10), Inches(0.5),
         "52G 협의체 — AI Native Product 거버넌스 수립을 위한 사례 발표", size=18, color=GRAY)
add_text(slide, Inches(1.5), Inches(4.2), Inches(10), Inches(0.5),
         "2026. 03. 12", size=16, color=GRAY)


# ============================================================
# SLIDE 2: Product Overview
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.6),
         "01  프로덕트 개요", size=30, color=ACCENT, bold=True)

# Divider
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.05), Inches(11.7), Pt(1))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0x33, 0x44, 0x55)
shape.line.fill.background()

# Description
add_text(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.8),
         "발전소 현장의 밸브·계기를 P&ID 도면 위에서 검색하고,\n360° 파노라마 로드뷰로 현장을 확인하며, EAM 정비이력까지 한 화면에서 조회하는 웹 애플리케이션",
         size=18, color=LIGHT)

# Feature cards
features = [
    ("P&ID 도면 뷰어", "도면 위 밸브 위치\n검색 및 마커 표시"),
    ("360° 로드뷰", "226개 씬 파노라마\n현장 실사 탐색"),
    ("EAM 정비이력", "HxGN EAM 연동\n실시간 이력 조회"),
    ("밸브/계기 사전", "카테고리별 분류\n격리 절차 안내"),
    ("3D 밸브 모델", "Three.js 기반\n입체 시각화"),
]

card_w = Inches(2.1)
card_h = Inches(2.2)
start_x = Inches(0.8)
gap = Inches(0.25)

for i, (title, desc) in enumerate(features):
    x = start_x + i * (card_w + gap)
    y = Inches(2.5)
    card = add_shape_bg(slide, x, y, card_w, card_h, BG_CARD)
    add_text(slide, x + Inches(0.2), y + Inches(0.3), card_w - Inches(0.4), Inches(0.5),
             title, size=18, color=ACCENT2, bold=True)
    add_text(slide, x + Inches(0.2), y + Inches(0.9), card_w - Inches(0.4), Inches(1),
             desc, size=14, color=LIGHT)

# Bottom stats
stats = [("개발 기간", "약 2개월"), ("커밋 수", "48회"), ("시작일", "2026.01.21"), ("개발 도구", "Claude Code")]
stat_w = Inches(2.6)
for i, (label, val) in enumerate(stats):
    x = Inches(0.8) + i * (stat_w + Inches(0.2))
    y = Inches(5.2)
    add_text(slide, x, y, stat_w, Inches(0.35), label, size=13, color=GRAY)
    add_text(slide, x, y + Inches(0.35), stat_w, Inches(0.4), val, size=20, color=WHITE, bold=True)


# ============================================================
# SLIDE 3: Architecture
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.6),
         "02  아키텍처", size=30, color=ACCENT, bold=True)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.05), Inches(11.7), Pt(1))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0x33, 0x44, 0x55)
shape.line.fill.background()

# 3-tier architecture
tiers = [
    ("Frontend", "Next.js 15 · React 19 · Tailwind CSS · Three.js", Inches(1.4), ACCENT, [
        "P&ID 도면 뷰어 (Canvas)", "360° 파노라마 로드뷰", "3D 밸브 모델 뷰어", "밸브 상세 패널 / 사전"
    ]),
    ("BFF (API Routes)", "Next.js Server-side API Routes", Inches(3.3), ACCENT2, [
        "/api/maintenance — EAM 프록시", "/api/diaries — 인계일지 프록시", "/api/roadview-settings — 설정 관리"
    ]),
    ("External / Legacy", "사내 시스템 및 클라우드 서비스", Inches(5.0), ORANGE, [
        "HxGN EAM REST API (정비이력)", "AWS Lambda (인계일지)", "JSON 파일 (밸브 마스터 데이터)"
    ]),
]

for title, subtitle, y, color, items in tiers:
    # Tier box
    box = add_shape_bg(slide, Inches(0.8), y, Inches(11.7), Inches(1.5), BG_CARD)

    # Color bar on left
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), y, Pt(5), Inches(1.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

    add_text(slide, Inches(1.2), y + Inches(0.1), Inches(3), Inches(0.4),
             title, size=20, color=color, bold=True)
    add_text(slide, Inches(1.2), y + Inches(0.5), Inches(3.5), Inches(0.3),
             subtitle, size=12, color=GRAY)

    # Items on right
    for j, item in enumerate(items):
        col = j // 2
        row = j % 2
        add_text(slide, Inches(5.5) + col * Inches(3.5), y + Inches(0.2) + row * Inches(0.55), Inches(3.5), Inches(0.5),
                 f"▸  {item}", size=13, color=LIGHT)


# ============================================================
# SLIDE 4: AI Code Generation
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
         "03  AI Code Generation 활용", size=30, color=ACCENT, bold=True)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.05), Inches(11.7), Pt(1))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0x33, 0x44, 0x55)
shape.line.fill.background()

# Key message
add_text(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.6),
         "처음부터 끝까지 Claude Code(AI 코드 생성)로 개발 — 전문 개발자 없이 현업 담당자가 직접 구현",
         size=18, color=YELLOW, bold=True)

# Usage cards
ai_items = [
    ("프로젝트 구조 설계", "Next.js 스캐폴딩,\n컴포넌트 설계를\nAI와 대화하며 구성", "초기 구조"),
    ("핵심 기능 구현", "P&ID 뷰어, 360° 파노라마,\n3D 모델 — 전체 코드\nAI가 생성", "코어 개발"),
    ("레거시 API 연동", "EAM REST API 파싱,\nAWS Lambda 호출 코드를\n프롬프트로 작성", "API 연동"),
    ("UI/UX 디자인", "Tailwind 반응형 UI,\n카드 효과, 테마를\nAI가 디자인+코드 동시 생성", "디자인"),
    ("디버깅·리팩토링", "파노라마 뒤틀림 버그 등\n증상 설명 → AI가\n수정 코드 제공", "유지보수"),
]

for i, (title, desc, tag) in enumerate(ai_items):
    x = Inches(0.8) + i * Inches(2.45)
    y = Inches(2.2)
    card = add_shape_bg(slide, x, y, Inches(2.25), Inches(2.8), BG_CARD)

    # Tag
    add_text(slide, x + Inches(0.15), y + Inches(0.15), Inches(1.8), Inches(0.3),
             tag, size=11, color=ACCENT, bold=True)
    # Title
    add_text(slide, x + Inches(0.15), y + Inches(0.55), Inches(1.95), Inches(0.5),
             title, size=16, color=WHITE, bold=True)
    # Desc
    add_text(slide, x + Inches(0.15), y + Inches(1.15), Inches(1.95), Inches(1.4),
             desc, size=13, color=LIGHT)

# Bottom highlight
box = add_shape_bg(slide, Inches(0.8), Inches(5.4), Inches(11.7), Inches(1.2), RGBColor(0x15, 0x2B, 0x44))
add_text(slide, Inches(1.2), Inches(5.55), Inches(11), Inches(0.4),
         "핵심 포인트", size=14, color=ACCENT, bold=True)
add_text(slide, Inches(1.2), Inches(5.95), Inches(11), Inches(0.5),
         "AI Code Generation 덕분에 아이디어-프로토타입-프로덕트 사이클이 극단적으로 짧아짐.\n그러나 빠르게 만든 것을 안전하게 운영하기 위한 체계는 아직 부재.",
         size=15, color=LIGHT)


# ============================================================
# SLIDE 5: Legacy 연결의 어려움
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
         "04  레거시 연결의 어려움", size=30, color=ACCENT, bold=True)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.05), Inches(11.7), Pt(1))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0x33, 0x44, 0x55)
shape.line.fill.background()

# Left: 실제 겪은 어려움
add_shape_bg(slide, Inches(0.8), Inches(1.4), Inches(5.6), Inches(5.2), BG_CARD)
add_text(slide, Inches(1.1), Inches(1.6), Inches(5), Inches(0.4),
         "실제 겪은 어려움", size=20, color=ORANGE, bold=True)

difficulties = [
    ("EAM API 인증 체계", "Basic Auth·테넌트·조직코드 등 인증 정보를\n코드에 연결하는 과정에서 사내 보안 정책과의\n정합성 확인이 어려움"),
    ("사내망 vs 외부망", "개발은 외부(Vercel) 환경, 레거시(EAM·PMS)는\n사내망 — 네트워크 접근 자체가 장벽"),
    ("데이터 포맷 불일치", "EAM API 응답 구조가 문서와 다른 경우 빈번.\n실제 응답을 보며 파싱 로직 수차례 수정"),
]
for i, (title, desc) in enumerate(difficulties):
    y = Inches(2.2) + i * Inches(1.4)
    add_text(slide, Inches(1.1), y, Inches(5.2), Inches(0.3), title, size=16, color=ORANGE, bold=True)
    add_text(slide, Inches(1.1), y + Inches(0.35), Inches(5.2), Inches(0.9), desc, size=13, color=LIGHT)

# Right: 걱정되는 부분
add_shape_bg(slide, Inches(6.7), Inches(1.4), Inches(5.8), Inches(5.2), BG_CARD)
add_text(slide, Inches(7.0), Inches(1.6), Inches(5), Inches(0.4),
         "걱정되는 부분", size=20, color=RED, bold=True)

worries = [
    ("보안성 검토 통과", "사전 보안성 검토 25개 항목 중 대부분 미충족.\nAI 생성 코드는 기능에 집중 → 보안 요건 별도 필요"),
    ("코드 품질 보증", "AI 생성 코드의 취약점 점검, 코드 리뷰\n프로세스가 부재한 상태"),
    ("운영 이관 기준 부재", "프로토타입 → 운영 시스템 전환 시\n어떤 수준까지 갖춰야 하는지 기준이 없음"),
]
for i, (title, desc) in enumerate(worries):
    y = Inches(2.2) + i * Inches(1.4)
    add_text(slide, Inches(7.0), y, Inches(5.2), Inches(0.3), title, size=16, color=RED, bold=True)
    add_text(slide, Inches(7.0), y + Inches(0.35), Inches(5.2), Inches(0.9), desc, size=13, color=LIGHT)


# ============================================================
# SLIDE 6: Governance
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
         "05  거버넌스 필요성", size=30, color=ACCENT, bold=True)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.05), Inches(11.7), Pt(1))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0x33, 0x44, 0x55)
shape.line.fill.background()

add_text(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.5),
         "AI로 빠르게 만들 수 있는 시대 — \"잘 만든 것\"을 \"안전하게 쓸 수 있게\" 하는 체계가 필요합니다",
         size=17, color=YELLOW, bold=True)

gov_items = [
    ("보안 기준", "AI 생성 코드의 사전 보안성 검토 기준\nAI 프로덕트 특성에 맞는 경량 체크리스트", ACCENT),
    ("코드 리뷰 정책", "AI가 생성한 코드를 누가,\n어떤 기준으로 리뷰할 것인지", ACCENT2),
    ("레거시 연동 가이드", "사내 API 접근 권한 발급 절차\n인증 방식 표준화", GREEN),
    ("배포·운영 이관 기준", "프로토타입 → 운영 전환 시\n최소 요구사항 (인증, 로깅, 모니터링)", ORANGE),
    ("데이터 거버넌스", "AI 학습에 활용 불가하도록\n사내 데이터 경계 설정", RED),
]

for i, (title, desc, color) in enumerate(gov_items):
    x = Inches(0.8) + (i % 3) * Inches(4.0)
    y = Inches(2.1) + (i // 3) * Inches(2.2)
    w = Inches(3.7)
    h = Inches(1.9)
    card = add_shape_bg(slide, x, y, w, h, BG_CARD)

    # Number circle
    num_shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.15), y + Inches(0.2), Inches(0.45), Inches(0.45))
    num_shape.fill.solid()
    num_shape.fill.fore_color.rgb = color
    num_shape.line.fill.background()
    tf = num_shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = str(i + 1)
    p.font.size = Pt(18)
    p.font.color.rgb = BG_DARK
    p.font.bold = True
    p.font.name = "맑은 고딕"
    p.alignment = PP_ALIGN.CENTER

    add_text(slide, x + Inches(0.7), y + Inches(0.2), w - Inches(0.9), Inches(0.4),
             title, size=18, color=color, bold=True)
    add_text(slide, x + Inches(0.7), y + Inches(0.7), w - Inches(0.9), Inches(1),
             desc, size=13, color=LIGHT)


# ============================================================
# SLIDE 7: Closing
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

# Accent line
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(3.4), Inches(13.333), Pt(3))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()

add_text(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(0.8),
         "\"잘 만든 것을, 안전하게 쓸 수 있게\"", size=36, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_text(slide, Inches(1.5), Inches(2.8), Inches(10), Inches(0.5),
         "올해 함께 거버넌스를 만들어 갑시다", size=22, color=ACCENT, bold=False, align=PP_ALIGN.CENTER)

add_text(slide, Inches(1.5), Inches(3.9), Inches(10), Inches(0.5),
         "GLANCE · AI Code Generation · 52G 협의체", size=16, color=GRAY, align=PP_ALIGN.CENTER)

add_text(slide, Inches(1.5), Inches(5.0), Inches(10), Inches(0.5),
         "감사합니다", size=28, color=WHITE, bold=True, align=PP_ALIGN.CENTER)


# Save
output_path = r"c:\Users\USER\Downloads\pid-viewer\GLANCE_AI_Governance_발표자료.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
