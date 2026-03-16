# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color palette
BG_DARK = RGBColor(0x0F, 0x17, 0x2A)
BG_CARD = RGBColor(0x1A, 0x25, 0x3C)
ACCENT = RGBColor(0x38, 0xBD, 0xF8)
ACCENT2 = RGBColor(0x22, 0xD3, 0xEE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x94, 0xA3, 0xB8)
LIGHT = RGBColor(0xCB, 0xD5, 0xE1)
ORANGE = RGBColor(0xFB, 0x92, 0x3C)
GREEN = RGBColor(0x4A, 0xDE, 0x80)
RED = RGBColor(0xF8, 0x71, 0x71)
YELLOW = RGBColor(0xFA, 0xCC, 0x15)
BG_HIGHLIGHT = RGBColor(0x15, 0x2B, 0x44)
PURPLE = RGBColor(0xA7, 0x8B, 0xFA)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    shape.adjustments[0] = 0.05
    return shape


def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "\ub9d1\uc740 \uace0\ub515"
    p.alignment = align
    return txBox


def add_multiline(slide, left, top, width, height, lines, size=16, color=LIGHT, line_spacing=1.3):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, txt_color, txt_bold, txt_size) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(txt_size if txt_size else size)
        run.font.color.rgb = txt_color if txt_color else color
        run.font.bold = txt_bold
        run.font.name = "\ub9d1\uc740 \uace0\ub515"
    return txBox


def add_divider(slide, y=Inches(1.05)):
    shape = add_rect(slide, Inches(0.8), y, Inches(11.7), Pt(1), RGBColor(0x33, 0x44, 0x55))
    return shape


def add_section_header(slide, number, title):
    add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
             f"{number}  {title}", size=30, color=ACCENT, bold=True)
    add_divider(slide)


def add_numbered_circle(slide, x, y, number, color):
    num_shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(0.45), Inches(0.45))
    num_shape.fill.solid()
    num_shape.fill.fore_color.rgb = color
    num_shape.line.fill.background()
    tf = num_shape.text_frame
    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.size = Pt(18)
    p.font.color.rgb = BG_DARK
    p.font.bold = True
    p.font.name = "\ub9d1\uc740 \uace0\ub515"
    p.alignment = PP_ALIGN.CENTER
    return num_shape


# ============================================================
# SLIDE 1: Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_rect(slide, Inches(0), Inches(3.1), Inches(13.333), Pt(3), ACCENT)
add_text(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(0.8),
         "AI \uac70\ubc84\ub10c\uc2a4 \ud544\uc694\uc131", size=48, color=WHITE, bold=True)
add_text(slide, Inches(1.5), Inches(2.4), Inches(10), Inches(0.7),
         "GLANCE \ud504\ub85c\uc81d\ud2b8 \uc0ac\ub840 \uae30\ubc18", size=28, color=ACCENT, bold=True)
add_text(slide, Inches(1.5), Inches(3.5), Inches(10), Inches(0.5),
         "52G \ud611\uc758\uccb4  |  AI Native Product \uac70\ubc84\ub10c\uc2a4 \uc218\ub9bd", size=18, color=GRAY)
add_text(slide, Inches(1.5), Inches(4.1), Inches(10), Inches(0.5),
         "2026. 03. 12", size=16, color=GRAY)


# ============================================================
# SLIDE 2: Why Governance?
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_section_header(slide, "01", "\uc65c \uac70\ubc84\ub10c\uc2a4\uac00 \ud544\uc694\ud55c\uac00?")

# Big message
add_shape_bg(slide, Inches(0.8), Inches(1.4), Inches(11.7), Inches(1.5), BG_HIGHLIGHT)
add_rect(slide, Inches(0.8), Inches(1.4), Pt(4), Inches(1.5), YELLOW)

add_text(slide, Inches(1.3), Inches(1.55), Inches(11), Inches(0.5),
         "AI Code Generation \uc2dc\ub300, \uc0c8\ub85c\uc6b4 \uac00\ub2a5\uc131\uc774 \uc5f4\ub838\uc2b5\ub2c8\ub2e4",
         size=22, color=WHITE, bold=True)
add_text(slide, Inches(1.3), Inches(2.1), Inches(11), Inches(0.6),
         "\ube44\uac1c\ubc1c\uc790(\uae30\uacc4 \uc5d4\uc9c0\ub2c8\uc5b4)\uac00 AI \ucf54\ub4dc \uc0dd\uc131 \ub3c4\uad6c\ub9cc\uc73c\ub85c \uc2e4\ubb34 \uc2dc\uc2a4\ud15c\uc744 \uad6c\ucd95\ud560 \uc218 \uc788\uac8c \ub418\uc5c8\uc2b5\ub2c8\ub2e4.\n\uc774\ub294 \ud601\uc2e0\uc758 \uae30\ud68c\uc774\uc9c0\ub9cc, \uccb4\uacc4\uc801\uc778 \uad00\ub9ac \uae30\uc900\uc774 \ud568\uaed8 \ud544\uc694\ud569\ub2c8\ub2e4.",
         size=16, color=LIGHT)

# Before / After comparison
add_text(slide, Inches(0.8), Inches(3.3), Inches(5.5), Inches(0.4),
         "\uac70\ubc84\ub10c\uc2a4 \ub3c4\uc785 \uc804", size=20, color=RED, bold=True)
add_text(slide, Inches(6.8), Inches(3.3), Inches(5.5), Inches(0.4),
         "\uac70\ubc84\ub10c\uc2a4 \ub3c4\uc785 \ud6c4", size=20, color=GREEN, bold=True)

before_items = [
    "\uac1c\uc778\uc774 \uc54c\uc544\uc11c \uac1c\ubc1c\xb7\ubc30\ud3ec",
    "\ubcf4\uc548 \uae30\uc900\uc744 \ucc98\uc74c\ubd80\ud130 \uc54c\uae30 \uc5b4\ub824\uc6c0",
    "\uadf8 \ub2e4\uc74c \ub2e8\uacc4\uc758 \uae38\uc774 \ubcf4\uc774\uc9c0 \uc54a\uc74c",
    "\ud488\uc9c8 \ud3b8\ucc28, \ubcf4\uc548 \uc0ac\uac01\uc9c0\ub300",
]

after_items = [
    "\ud68c\uc0ac\uac00 \ud568\uaed8 \uac00\uc774\ub4dc\xb7\uc9c0\uc6d0",
    "\uc0ac\uc804 \uccb4\ud06c\ub9ac\uc2a4\ud2b8\ub85c \ucd08\uae30\ubd80\ud130 \uc548\uc804\ud558\uac8c",
    "\ubc30\ud3ec\xb7\uc6b4\uc601 \ub85c\ub4dc\ub9f5 \uba85\ud655",
    "\uc548\uc804\ud55c \ud601\uc2e0, \uccb4\uacc4\uc801 \ud655\uc0b0",
]

# Before card
add_shape_bg(slide, Inches(0.8), Inches(3.8), Inches(5.5), Inches(3.0), BG_CARD)
add_rect(slide, Inches(0.8), Inches(3.8), Inches(5.5), Pt(3), RED)
for i, item in enumerate(before_items):
    add_text(slide, Inches(1.2), Inches(4.1) + i * Inches(0.6), Inches(4.8), Inches(0.5),
             f"\u2718  {item}", size=15, color=LIGHT)

# After card
add_shape_bg(slide, Inches(6.8), Inches(3.8), Inches(5.7), Inches(3.0), BG_CARD)
add_rect(slide, Inches(6.8), Inches(3.8), Inches(5.7), Pt(3), GREEN)
for i, item in enumerate(after_items):
    add_text(slide, Inches(7.2), Inches(4.1) + i * Inches(0.6), Inches(5.0), Inches(0.5),
             f"\u2714  {item}", size=15, color=GREEN)

# Arrow between
add_text(slide, Inches(6.0), Inches(4.8), Inches(1.0), Inches(0.6),
         "\u27a1", size=36, color=ACCENT, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 3: GLANCE Case - Current Status Table
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_section_header(slide, "02", "GLANCE \uc0ac\ub840 \u2014 \ud604\ud669\uacfc \uacfc\uc81c")

add_text(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.5),
         "GLANCE \ud504\ub85c\uc81d\ud2b8\uc5d0\uc11c \ub290\ub080 \uac70\ubc84\ub10c\uc2a4 \ud544\uc694 \uc601\uc5ed",
         size=17, color=YELLOW, bold=True)

# Table-like layout
table_data = [
    ("\ubcf4\uc548", "\ud658\uacbd\ubcc0\uc218\ub85c \ubd84\ub9ac \uc644\ub8cc", "\uc2dc\ud06c\ub9bf \ub9e4\ub2c8\uc800 \uc5f0\ub3d9\uc73c\ub85c \ucd94\uac00 \ubcf4\ud638", ACCENT, GREEN),
    ("\ubc30\ud3ec", "\uac1c\uc778 \uacc4\uc815(Vercel)\uc5d0 \ubc30\ud3ec \uc911", "\ud68c\uc0ac \uc2b9\uc778 \ubc30\ud3ec \ud658\uacbd \uc804\ud658", ACCENT2, YELLOW),
    ("\uc778\uc99d", "\ubcc4\ub3c4 \uc0ac\uc6a9\uc790 \uc778\uc99d \uc5c6\uc74c", "\uadf8\ub8f9\uc6e8\uc5b4(SSO) \uc5f0\ub3d9 \uac80\ud1a0", ORANGE, YELLOW),
    ("\ub370\uc774\ud130", "\ub3c4\uba74\xb7\ud30c\ub178\ub77c\ub9c8 \ub85c\ucee8 \uc800\uc7a5", "\uc624\ube0c\uc81d\ud2b8 \uc2a4\ud1a0\ub9ac\uc9c0(S3) \uc774\uad00", PURPLE, YELLOW),
    ("\ucf54\ub4dc \ud488\uc9c8", "AI \uc0dd\uc131 \ucf54\ub4dc \uac80\uc99d \ubbf8\ube44", "\ucf54\ub4dc \ub9ac\ubdf0 \ubc0f \uac80\uc218 \ud504\ub85c\uc138\uc2a4 \ub9c8\ub828", RED, YELLOW),
]

# Headers
add_shape_bg(slide, Inches(0.8), Inches(1.9), Inches(2.6), Inches(0.55), ACCENT)
add_text(slide, Inches(0.9), Inches(1.98), Inches(2.4), Inches(0.4),
         "\uc601\uc5ed", size=16, color=BG_DARK, bold=True, align=PP_ALIGN.CENTER)

add_shape_bg(slide, Inches(3.5), Inches(1.9), Inches(4.3), Inches(0.55), ACCENT)
add_text(slide, Inches(3.6), Inches(1.98), Inches(4.1), Inches(0.4),
         "\ud604\uc7ac \uc0c1\ud0dc", size=16, color=BG_DARK, bold=True, align=PP_ALIGN.CENTER)

add_shape_bg(slide, Inches(7.9), Inches(1.9), Inches(4.6), Inches(0.55), ACCENT)
add_text(slide, Inches(8.0), Inches(1.98), Inches(4.4), Inches(0.4),
         "\uac70\ubc84\ub10c\uc2a4\ub85c \ud574\uacb0", size=16, color=BG_DARK, bold=True, align=PP_ALIGN.CENTER)

for i, (area, current, governance, area_color, status_color) in enumerate(table_data):
    y = Inches(2.6) + i * Inches(0.85)

    # Area column
    add_shape_bg(slide, Inches(0.8), y, Inches(2.6), Inches(0.7), BG_CARD)
    add_rect(slide, Inches(0.8), y, Pt(4), Inches(0.7), area_color)
    add_text(slide, Inches(1.1), y + Inches(0.15), Inches(2.2), Inches(0.4),
             area, size=16, color=area_color, bold=True)

    # Current status column
    add_shape_bg(slide, Inches(3.5), y, Inches(4.3), Inches(0.7), BG_CARD)
    add_text(slide, Inches(3.7), y + Inches(0.15), Inches(3.9), Inches(0.4),
             current, size=14, color=LIGHT)

    # Governance solution column
    add_shape_bg(slide, Inches(7.9), y, Inches(4.6), Inches(0.7), BG_CARD)
    add_text(slide, Inches(8.1), y + Inches(0.15), Inches(4.2), Inches(0.4),
             governance, size=14, color=status_color)


# ============================================================
# SLIDE 4: Governance Direction - 3 Pillars
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_section_header(slide, "03", "\uac70\ubc84\ub10c\uc2a4\uc758 \ubc29\ud5a5 \u2014 \ub9c9\ub294 \uac83\uc774 \uc544\ub2c8\ub77c \ub3d5\ub294 \uac83")

# Subtitle
add_shape_bg(slide, Inches(0.8), Inches(1.4), Inches(11.7), Inches(1.0), BG_HIGHLIGHT)
add_rect(slide, Inches(0.8), Inches(1.4), Pt(4), Inches(1.0), YELLOW)
add_text(slide, Inches(1.3), Inches(1.55), Inches(11), Inches(0.7),
         "AI \ucf54\ub4dc \uc0dd\uc131\uc740 \ud604\uc7a5\uc758 \ubb38\uc81c\ub97c \ud604\uc7a5 \ub2f4\ub2f9\uc790\uac00 \uc9c1\uc811 \ud574\uacb0\ud560 \uc218 \uc788\uac8c \ud569\ub2c8\ub2e4.\n\uac70\ubc84\ub10c\uc2a4\ub294 \uc774 \ud601\uc2e0\uc744 \ub9c9\ub294 \uc7a5\ubcbd\uc774 \uc544\ub2c8\ub77c, \uc548\uc804\ud558\uac8c \ud655\uc0b0\uc2dc\ud0a4\ub294 \uac00\ub4dc\ub808\uc77c\uc785\ub2c8\ub2e4.",
         size=16, color=LIGHT)

# 3 big pillars
pillars = [
    ("01", "\uac00\uc774\ub4dc\ub77c\uc778 \uc81c\uacf5", ACCENT,
     [
         "AI \ucf54\ub4dc \uc0dd\uc131 \uc2dc \uc9c0\ucf1c\uc57c \ud560",
         "\ubcf4\uc548 \uae30\ubcf8 \uc6d0\uce59",
         "(HTTPS, \uc778\uc99d\uc815\ubcf4 \ubd84\ub9ac, \ub370\uc774\ud130 \ubd84\ub958)",
         "",
         "\ube44\uac1c\ubc1c\uc790\ub3c4 \ub530\ub77c\ud560 \uc218 \uc788\ub294",
         "\uccb4\ud06c\ub9ac\uc2a4\ud2b8 \ud615\ud0dc",
     ]),
    ("02", "\ubc30\ud3ec\xb7\uc6b4\uc601 \ud658\uacbd \uc9c0\uc6d0", ACCENT2,
     [
         "\ud68c\uc0ac \ucc28\uc6d0\uc758 \ubc30\ud3ec \ud658\uacbd \uc81c\uacf5",
         "(\uc0ac\ub0b4 \uc11c\ubc84, \ud68c\uc0ac \ud074\ub77c\uc6b0\ub4dc)",
         "",
         "\uc778\uc99d\uc815\ubcf4 \uad00\ub9ac \uc778\ud504\ub77c \uc9c0\uc6d0",
         "(Secret Manager, \ud658\uacbd\ubcc0\uc218 \uad00\ub9ac)",
         "",
         "\ub300\uc6a9\ub7c9 \ud30c\uc77c \uc800\uc7a5\uc18c \uc81c\uacf5 (S3)",
     ]),
    ("03", "\ud655\uc0b0 \uccb4\uacc4 \uad6c\ucd95", ORANGE,
     [
         "GLANCE \uac19\uc740 \uc6b0\uc218 \uc0ac\ub840\ub97c",
         "\ub2e4\ub978 \ubc1c\uc804\uc18c\xb7\ubd80\uc11c\ub85c \ud655\uc0b0",
         "",
         "AI \ud65c\uc6a9 \uad50\uc721 \ud504\ub85c\uadf8\ub7a8",
         "",
         "\ud504\ub85c\ud1a0\ud0c0\uc785 \u2192 \uc6b4\uc601 \uc804\ud658 \uae30\uc900",
         "\uba85\ud655\ud654",
     ]),
]

for i, (num, title, color, items) in enumerate(pillars):
    x = Inches(0.8) + i * Inches(4.0)
    y = Inches(2.8)
    w = Inches(3.7)
    h = Inches(4.2)

    add_shape_bg(slide, x, y, w, h, BG_CARD)
    add_rect(slide, x, y, w, Pt(4), color)

    # Number + title
    add_numbered_circle(slide, x + Inches(0.2), y + Inches(0.25), num, color)
    add_text(slide, x + Inches(0.75), y + Inches(0.25), w - Inches(1.0), Inches(0.4),
             title, size=22, color=color, bold=True)

    # Items
    for j, item in enumerate(items):
        add_text(slide, x + Inches(0.3), y + Inches(0.9) + j * Inches(0.4), w - Inches(0.6), Inches(0.4),
                 item, size=13, color=LIGHT)


# ============================================================
# SLIDE 5: 5 Key Areas
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_section_header(slide, "04", "\uac70\ubc84\ub10c\uc2a4 5\ub300 \uc601\uc5ed")

areas = [
    ("\ubcf4\uc548 \uae30\uc900", ACCENT,
     "\uc778\uc99d\uc815\ubcf4 \ud658\uacbd\ubcc0\uc218 \ubd84\ub9ac\n\uc2dc\ud06c\ub9bf \ub9e4\ub2c8\uc800 \uc5f0\ub3d9\nHTTPS \ud1b5\uc2e0 \ud544\uc218\ud654\n\uc0ac\uc804 \ubcf4\uc548\uc131 \uac80\ud1a0 \uccb4\ud06c\ub9ac\uc2a4\ud2b8"),
    ("\ubc30\ud3ec \ud658\uacbd", ACCENT2,
     "\ud68c\uc0ac \uc2b9\uc778 \ubc30\ud3ec \ud658\uacbd \uc81c\uacf5\n\uc0ac\ub0b4 \uc11c\ubc84 \ub610\ub294 \ud68c\uc0ac \ud074\ub77c\uc6b0\ub4dc\nIP \uc811\uadfc\uc81c\uc5b4 \ubc0f \uc811\uadfc \uad00\ub9ac"),
    ("\ub370\uc774\ud130 \uad00\ub9ac", ORANGE,
     "\ub3c4\uba74, \ud30c\ub178\ub77c\ub9c8 \ub4f1 \ub300\uc6a9\ub7c9 \ud30c\uc77c\n\uc624\ube0c\uc81d\ud2b8 \uc2a4\ud1a0\ub9ac\uc9c0(S3) \uc774\uad00\n\ub370\uc774\ud130 \ubd84\ub958 \ubc0f \uc811\uadfc \uad8c\ud55c \uae30\uc900"),
    ("\ucf54\ub4dc \ud488\uc9c8", PURPLE,
     "AI \uc0dd\uc131 \ucf54\ub4dc \ub9ac\ubdf0 \ud504\ub85c\uc138\uc2a4\n\ud488\uc9c8 \uac80\uc218 \uae30\uc900 \ub9c8\ub828\n\ucf54\ub4dc \ubcf4\uc548 \uc810\uac80 \ub3c4\uad6c \ud65c\uc6a9"),
    ("\ud655\uc0b0 \uccb4\uacc4", GREEN,
     "\uc6b0\uc218 \uc0ac\ub840 \uacf5\uc720 \ubc0f \uad50\uc721\nAI \ud65c\uc6a9 \ud604\uc5c5 \ucee4\ubba4\ub2c8\ud2f0\n\ud504\ub85c\ud1a0\ud0c0\uc785 \u2192 \uc6b4\uc601 \uc804\ud658 \uac00\uc774\ub4dc"),
]

for i, (title, color, desc) in enumerate(areas):
    # 5 cards: 3 on top, 2 on bottom centered
    if i < 3:
        x = Inches(0.8) + i * Inches(4.0)
        y = Inches(1.4)
    else:
        x = Inches(2.8) + (i - 3) * Inches(4.0)
        y = Inches(4.2)
    w = Inches(3.7)
    h = Inches(2.5)

    add_shape_bg(slide, x, y, w, h, BG_CARD)
    add_rect(slide, x, y, w, Pt(4), color)

    add_numbered_circle(slide, x + Inches(0.2), y + Inches(0.25), i + 1, color)
    add_text(slide, x + Inches(0.75), y + Inches(0.25), w - Inches(1.0), Inches(0.4),
             title, size=20, color=color, bold=True)
    add_text(slide, x + Inches(0.3), y + Inches(0.85), w - Inches(0.6), Inches(1.5),
             desc, size=13, color=LIGHT)


# ============================================================
# SLIDE 6: Expected Outcomes
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_section_header(slide, "05", "\uae30\ub300 \ud6a8\uacfc")

# Big flow
flow_items = [
    ("\ud604\uc7a5\uc758 \uc544\uc774\ub514\uc5b4", GRAY, "\ub2f4\ub2f9\uc790\uac00 \ubb38\uc81c\ub97c \ubc1c\uacac"),
    ("AI\ub85c \ube60\ub978 \uad6c\ud604", ACCENT, "\ucf54\ub4dc \uc0dd\uc131\uc73c\ub85c \ud504\ub85c\ud1a0\ud0c0\uc785"),
    ("\uac70\ubc84\ub10c\uc2a4 \uac80\uc99d", YELLOW, "\ubcf4\uc548\xb7\ubc30\ud3ec\xb7\ud488\uc9c8 \uccb4\ud06c"),
    ("\uc548\uc804\ud55c \uc6b4\uc601", GREEN, "\ud604\uc5c5\uc5d0\uc11c \uc2e4\uc81c \uc0ac\uc6a9"),
    ("\ub2e4\ub978 \ud604\uc7a5\uc73c\ub85c \ud655\uc0b0", ORANGE, "\uc131\uacf5 \uc0ac\ub840 \uacf5\uc720"),
]

for i, (title, color, desc) in enumerate(flow_items):
    x = Inches(0.5) + i * Inches(2.5)
    y = Inches(1.6)

    # Circle with number
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.75), y, Inches(0.9), Inches(0.9))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = str(i + 1)
    p.font.size = Pt(28)
    p.font.color.rgb = BG_DARK
    p.font.bold = True
    p.font.name = "\ub9d1\uc740 \uace0\ub515"

    add_text(slide, x, y + Inches(1.1), Inches(2.4), Inches(0.5),
             title, size=16, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x, y + Inches(1.55), Inches(2.4), Inches(0.4),
             desc, size=12, color=LIGHT, align=PP_ALIGN.CENTER)

    # Arrow between
    if i < len(flow_items) - 1:
        add_text(slide, x + Inches(2.2), y + Inches(0.15), Inches(0.5), Inches(0.6),
                 "\u27a1", size=28, color=GRAY, align=PP_ALIGN.CENTER)

# 3 benefit cards
benefits = [
    ("\ub9ac\uc2a4\ud06c \ucd5c\uc18c\ud654", ACCENT,
     "AI\ub85c \ube60\ub974\uac8c \ub9cc\ub4e4\ub418,\n\ubcf4\uc548\xb7\ud488\uc9c8\uc740 \uccb4\uacc4\uc801\uc73c\ub85c \uad00\ub9ac\n\u2192 \uc548\uc804\ud55c \ud601\uc2e0"),
    ("\ud6a8\uacfc \uadf9\ub300\ud654", ACCENT2,
     "\ud604\uc5c5 \ub2f4\ub2f9\uc790\uc758 \ub3c4\uba54\uc778 \uc9c0\uc2dd\uc774\n\uc9c1\uc811 \uc2dc\uc2a4\ud15c\uc73c\ub85c \uc5f0\uacb0\n\u2192 \uc2e4\uc9c8\uc801 \uc5c5\ubb34 \ud6a8\uc728 \ud5a5\uc0c1"),
    ("52G \uc804\uccb4 \ud655\uc0b0", GREEN,
     "\uc131\uacf5 \uc0ac\ub840\ub97c \uae30\ubc18\uc73c\ub85c\n\ub2e4\ub978 \ubc1c\uc804\uc18c\xb7\ubd80\uc11c\ub85c \ud655\ub300\n\u2192 AI \ud65c\uc6a9 \ubb38\ud654 \uc815\ucc29"),
]

for i, (title, color, desc) in enumerate(benefits):
    x = Inches(0.8) + i * Inches(4.0)
    y = Inches(3.8)
    w = Inches(3.7)
    h = Inches(2.6)

    add_shape_bg(slide, x, y, w, h, BG_CARD)
    add_rect(slide, x, y, w, Pt(4), color)
    add_text(slide, x + Inches(0.3), y + Inches(0.25), w - Inches(0.6), Inches(0.4),
             title, size=20, color=color, bold=True)
    add_text(slide, x + Inches(0.3), y + Inches(0.85), w - Inches(0.6), Inches(1.5),
             desc, size=14, color=LIGHT)


# ============================================================
# SLIDE 7: MISO vs Claude Code
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_section_header(slide, "06", "MISO(Defy) vs Claude Code \u2014 \uacbd\uc7c1\uc774 \uc544\ub2c8\ub77c \ubcf4\uc644")

# Subtitle
add_shape_bg(slide, Inches(0.8), Inches(1.4), Inches(11.7), Inches(0.8), BG_HIGHLIGHT)
add_rect(slide, Inches(0.8), Inches(1.4), Pt(4), Inches(0.8), YELLOW)
add_text(slide, Inches(1.3), Inches(1.5), Inches(11), Inches(0.5),
         "\ub450 AI\ub294 \uc6a9\ub3c4\uac00 \ub2e4\ub985\ub2c8\ub2e4. \uac01\uc790 \uc798\ud558\ub294 \uc601\uc5ed\uc774 \ub2e4\ub978 \ubcf4\uc644 \uad00\uacc4\uc785\ub2c8\ub2e4.",
         size=17, color=LIGHT)

# MISO Card (Left)
add_shape_bg(slide, Inches(0.8), Inches(2.5), Inches(5.5), Inches(4.3), BG_CARD)
add_rect(slide, Inches(0.8), Inches(2.5), Inches(5.5), Pt(4), PURPLE)

add_text(slide, Inches(1.2), Inches(2.7), Inches(4.8), Inches(0.5),
         "MISO (Defy)", size=28, color=PURPLE, bold=True)
add_text(slide, Inches(1.2), Inches(3.3), Inches(4.8), Inches(0.4),
         "\ub611\ub611\ud55c \ube44\uc11c", size=18, color=YELLOW, bold=True)

miso_items = [
    ("\uc6a9\ub3c4", "\uc5c5\ubb34 \uc9c0\uc6d0 \ucc57\ubd07 (\uc9c8\uc758\uc751\ub2f5, \ubb38\uc11c \uc694\uc57d)"),
    ("\uacb0\uacfc\ubb3c", "\ud14d\uc2a4\ud2b8 \ub2f5\ubcc0"),
    ("\uc608\uc2dc", "\"\ubc38\ube0c \uc815\ube44\uc774\ub825 \uc54c\ub824\uc918\" \u2192 \ud14d\uc2a4\ud2b8\ub85c \ub2f5\ubcc0"),
    ("\uac15\uc810", "\uc0ac\ub0b4 \ub370\uc774\ud130 \uae30\ubc18 \ube60\ub978 \uc751\ub2f5"),
]

for i, (label, desc) in enumerate(miso_items):
    y = Inches(3.85) + i * Inches(0.65)
    add_text(slide, Inches(1.2), y, Inches(1.3), Inches(0.4),
             label, size=14, color=PURPLE, bold=True)
    add_text(slide, Inches(2.6), y, Inches(3.4), Inches(0.4),
             desc, size=14, color=LIGHT)

# Claude Code Card (Right)
add_shape_bg(slide, Inches(7.0), Inches(2.5), Inches(5.5), Inches(4.3), BG_CARD)
add_rect(slide, Inches(7.0), Inches(2.5), Inches(5.5), Pt(4), ACCENT)

add_text(slide, Inches(7.4), Inches(2.7), Inches(4.8), Inches(0.5),
         "Claude Code", size=28, color=ACCENT, bold=True)
add_text(slide, Inches(7.4), Inches(3.3), Inches(4.8), Inches(0.4),
         "\uac1c\ubc1c\uc790", size=18, color=YELLOW, bold=True)

claude_items = [
    ("\uc6a9\ub3c4", "\uc18c\ud504\ud2b8\uc6e8\uc5b4 \uac1c\ubc1c (\ucf54\ub4dc \uc0dd\uc131)"),
    ("\uacb0\uacfc\ubb3c", "\uc2e4\uc81c \ub3d9\uc791\ud558\ub294 \ud504\ub85c\ub355\ud2b8"),
    ("\uc608\uc2dc", "\"\ubc38\ube0c \uc870\ud68c \uc2dc\uc2a4\ud15c \ub9cc\ub4e4\uc5b4\uc918\" \u2192 GLANCE \ud0c4\uc0dd"),
    ("\uac15\uc810", "\ube44\uac1c\ubc1c\uc790\ub3c4 \uc2dc\uc2a4\ud15c \uad6c\ucd95 \uac00\ub2a5"),
]

for i, (label, desc) in enumerate(claude_items):
    y = Inches(3.85) + i * Inches(0.65)
    add_text(slide, Inches(7.4), y, Inches(1.3), Inches(0.4),
             label, size=14, color=ACCENT, bold=True)
    add_text(slide, Inches(8.8), y, Inches(3.4), Inches(0.4),
             desc, size=14, color=LIGHT)

# VS in the middle
vs_shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.95), Inches(4.0), Inches(1.1), Inches(1.1))
vs_shape.fill.solid()
vs_shape.fill.fore_color.rgb = BG_DARK
vs_shape.line.color.rgb = YELLOW
vs_shape.line.width = Pt(2)
tf = vs_shape.text_frame
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
p = tf.paragraphs[0]
p.text = "VS"
p.font.size = Pt(24)
p.font.color.rgb = YELLOW
p.font.bold = True
p.font.name = "\ub9d1\uc740 \uace0\ub515"

# Bottom conclusion
add_shape_bg(slide, Inches(2.5), Inches(6.3), Inches(8.3), Inches(0.7), BG_HIGHLIGHT)
add_rect(slide, Inches(2.5), Inches(6.3), Pt(4), Inches(0.7), GREEN)
add_text(slide, Inches(2.8), Inches(6.38), Inches(7.8), Inches(0.5),
         "MISO\ub85c \ubb3c\uc5b4\ubcf4\uace0, Claude Code\ub85c \ub9cc\ub4e0\ub2e4 \u2014 \ud568\uaed8 \uc4f0\uba74 \ub354 \uac15\ub825\ud569\ub2c8\ub2e4",
         size=17, color=GREEN, bold=True, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 8: Closing
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_rect(slide, Inches(0), Inches(3.6), Inches(13.333), Pt(3), ACCENT)

add_text(slide, Inches(1.5), Inches(1.3), Inches(10), Inches(0.7),
         "\uac70\ubc84\ub10c\uc2a4\ub294", size=26, color=GRAY, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(1.0),
         "\ud601\uc2e0\uc744 \ub9c9\ub294 \uac83\uc774 \uc544\ub2c8\ub77c,\n\ud601\uc2e0\uc774 \uc0b4\uc544\ub0a8\ub294 \uccb4\uacc4\ub97c \ub9cc\ub4dc\ub294 \uac83\uc785\ub2c8\ub2e4.",
         size=32, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_text(slide, Inches(1.5), Inches(4.1), Inches(10), Inches(0.7),
         "\ube60\ub974\uac8c \ub9cc\ub4e4 \uc218 \uc788\ub294 \uc774 \ud798\uc744,\n\uc870\uc9c1 \uc548\uc5d0\uc11c \uc548\uc804\ud558\uac8c \ud65c\uc6a9\ud560 \uc218 \uc788\ub294 \uccb4\uacc4\ub97c \ud568\uaed8 \ub9cc\ub4e4\uc5b4\uac11\uc2dc\ub2e4.",
         size=20, color=ACCENT, align=PP_ALIGN.CENTER)

add_text(slide, Inches(1.5), Inches(5.5), Inches(10), Inches(0.5),
         "\uac10\uc0ac\ud569\ub2c8\ub2e4", size=28, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_text(slide, Inches(1.5), Inches(6.2), Inches(10), Inches(0.4),
         "GLANCE  \xb7  AI Governance  \xb7  52G \ud611\uc758\uccb4", size=14, color=GRAY, align=PP_ALIGN.CENTER)


# Save
output_path = r"c:\Users\USER\Downloads\pid-viewer\GLANCE_Governance_v2.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
