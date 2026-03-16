"""
GLANCE 발표자료 마지막에 '현장에 답이 있다' 마무리 장표 추가
기존 스타일: 배경 #0F172A(다크네이비), 액센트 #38BDF8(시안), 폰트 맑은 고딕
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import copy

# 기존 파일 열기
prs = Presentation('c:/Users/USER/Downloads/pid-viewer/GLANCE_AI_Governance_v2.pptx')

# 색상 정의
BG_COLOR = RGBColor(0x0F, 0x17, 0x2A)       # 다크 네이비
CYAN = RGBColor(0x38, 0xBD, 0xF8)            # 시안 액센트
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCB, 0xD5, 0xE1)      # 연한 회색
DIM_GRAY = RGBColor(0x64, 0x74, 0x8B)        # 어두운 회색
WARM_AMBER = RGBColor(0xFB, 0xBF, 0x24)      # 따뜻한 앰버/골드
FONT_NAME = "맑은 고딕"

slide_width = prs.slide_width   # 12192000
slide_height = prs.slide_height  # 6858000

def add_text_box(slide, left, top, width, height, text, font_size, font_color=WHITE,
                 bold=False, alignment=PP_ALIGN.CENTER, font_name=FONT_NAME, line_spacing=None):
    from pptx.util import Pt
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    if line_spacing:
        p.line_spacing = Pt(line_spacing)
    run = p.runs[0]
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.color.rgb = font_color
    run.font.bold = bold
    return txBox

def add_multiline_box(slide, left, top, width, height, lines, alignment=PP_ALIGN.CENTER):
    """여러 줄을 하나의 텍스트박스에 추가 (각 줄마다 스타일 지정)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, (text, font_size, font_color, bold) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = alignment
        run = p.add_run()
        run.text = text
        run.font.name = FONT_NAME
        run.font.size = Pt(font_size)
        run.font.color.rgb = font_color
        run.font.bold = bold
    return txBox

def add_rect(slide, left, top, width, height, fill_color):
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

# ============================================================
# 슬라이드 10: "현장에 답이 있다" — 임팩트 마무리
# ============================================================
slide_layout = prs.slide_layouts[6]  # blank
slide = prs.slides.add_slide(slide_layout)

# 배경색 설정
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = BG_COLOR

# --- 상단: 시안 액센트 라인 (기존 스타일 유지) ---
add_rect(slide, 0, Emu(0), slide_width, Emu(76200), CYAN)

# --- 메인 콘텐츠 ---

# 작은 시안 태그라인
add_text_box(slide, Emu(1371600), Emu(1143000), Emu(9144000), Emu(365760),
             "GLANCE가 시작된 이유", 16, CYAN, bold=False)

# 메인 카피 1: "현장에 답이 있다"
add_text_box(slide, Emu(1371600), Emu(1600200), Emu(9144000), Emu(914400),
             "현장에 답이 있다", 44, WHITE, bold=True)

# 시안 구분선 (짧은 가로선)
add_rect(slide, Emu(5400000), Emu(2600000), Emu(1200000), Emu(38100), CYAN)

# 본문 메시지 블록
lines_main = [
    ("도면 한 장 펼쳐놓고 밸브 하나 찾는 데 10분,", 18, LIGHT_GRAY, False),
    ("정비이력 확인하려면 시스템 세 개를 돌아다녀야 했습니다.", 18, LIGHT_GRAY, False),
    ("", 10, LIGHT_GRAY, False),
    ("그래서 직접 만들었습니다.", 22, WHITE, True),
    ("현장의 불편함을 가장 잘 아는 사람이,", 18, LIGHT_GRAY, False),
    ("AI와 함께 직접 해결한 이야기입니다.", 18, LIGHT_GRAY, False),
]
add_multiline_box(slide, Emu(1371600), Emu(2800000), Emu(9144000), Emu(2000000), lines_main)

# 하단 강조 메시지 — 골드 컬러로 임팩트
add_text_box(slide, Emu(1371600), Emu(5000000), Emu(9144000), Emu(640080),
             "기술이 현장을 바꾸는 것이 아니라,", 20, LIGHT_GRAY, False)

add_text_box(slide, Emu(1371600), Emu(5400000), Emu(9144000), Emu(640080),
             "현장이 기술의 방향을 정합니다.", 26, WARM_AMBER, bold=True)

# 하단 시안 구분선
add_rect(slide, Emu(0), Emu(6300000), slide_width, Emu(38100), CYAN)

# 하단 푸터
add_text_box(slide, Emu(1371600), Emu(6400000), Emu(9144000), Emu(365760),
             "GLANCE · 현장에서 시작된 AI 프로덕트", 12, DIM_GRAY, bold=False)


# ============================================================
# 슬라이드 11: 보너스 — "감사합니다" + 핵심 한 줄
# ============================================================
slide2 = prs.slides.add_slide(slide_layout)

# 배경
bg2 = slide2.background
fill2 = bg2.fill
fill2.solid()
fill2.fore_color.rgb = BG_COLOR

# 상단 시안 라인
add_rect(slide2, 0, Emu(0), slide_width, Emu(76200), CYAN)

# 가운데 큰 "감사합니다"
add_text_box(slide2, Emu(1371600), Emu(1800000), Emu(9144000), Emu(914400),
             "감사합니다", 48, WHITE, bold=True)

# 구분선
add_rect(slide2, Emu(5000000), Emu(2900000), Emu(2000000), Emu(38100), CYAN)

# 핵심 메시지
lines_closing = [
    ("코딩을 몰라도, 현장을 알면", 22, LIGHT_GRAY, False),
    ("프로덕트를 만들 수 있는 시대.", 22, LIGHT_GRAY, False),
    ("", 10, LIGHT_GRAY, False),
    ("현장의 목소리가 곧 기술의 출발점입니다.", 24, CYAN, True),
]
add_multiline_box(slide2, Emu(1371600), Emu(3200000), Emu(9144000), Emu(1600000), lines_closing)

# 하단 키워드
add_text_box(slide2, Emu(1371600), Emu(5200000), Emu(9144000), Emu(457200),
             "GLANCE  ·  AI Code Generation  ·  현장에 답이 있다", 16, DIM_GRAY, bold=False)

# 하단 시안 라인
add_rect(slide2, Emu(0), Emu(6300000), slide_width, Emu(38100), CYAN)

# 이름/소속 (필요 시)
add_text_box(slide2, Emu(1371600), Emu(5700000), Emu(9144000), Emu(365760),
             "인터코 GON · 정승환", 14, LIGHT_GRAY, bold=False)

# 저장
output_path = 'c:/Users/USER/Downloads/pid-viewer/GLANCE_AI_Governance_v2.pptx'
prs.save(output_path)
print(f"저장 완료: {output_path}")
print(f"총 슬라이드 수: {len(prs.slides)}")
