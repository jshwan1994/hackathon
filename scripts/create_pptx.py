"""
GLANCE - P&ID Viewer 발표자료 생성 스크립트
고도화 버전 (2026.03 업데이트)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── 색상 정의 ──
DARK_BG = RGBColor(0x0F, 0x17, 0x2A)       # 진한 남색 배경
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)  # 주황 강조
ACCENT_BLUE = RGBColor(0x3B, 0x82, 0xF6)    # 파란 강조
ACCENT_GREEN = RGBColor(0x10, 0xB9, 0x81)   # 초록 강조
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MEDIUM_GRAY = RGBColor(0x64, 0x64, 0x64)
LIGHT_BG = RGBColor(0xF8, 0xF8, 0xF8)       # 밝은 배경
TAG_ORANGE_BG = RGBColor(0xFF, 0x6B, 0x2B)   # 태그 배경 주황

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


def add_dark_bg(slide):
    """어두운 배경 + 사각형 덮기"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG


def add_light_bg(slide):
    """밝은 배경"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = LIGHT_BG


def add_tag(slide, text, left, top, color=TAG_ORANGE_BG, width=Inches(2.2), height=Inches(0.55)):
    """좌상단 태그 라벨"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE


def add_title_text(slide, text, left, top, width, height, size=Pt(36), color=DARK_GRAY, bold=True, alignment=PP_ALIGN.LEFT):
    """제목 텍스트"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = size
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return tf


def add_body_text(slide, text, left, top, width, height, size=Pt(18), color=DARK_GRAY, bold=False, alignment=PP_ALIGN.LEFT):
    """본문 텍스트"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = size
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return tf


def add_multiline_text(slide, lines, left, top, width, height, size=Pt(16), color=DARK_GRAY, line_spacing=Pt(28)):
    """여러 줄 텍스트"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, is_bold, text_color, text_size) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = text_size or size
        p.font.bold = is_bold
        p.font.color.rgb = text_color or color
        p.space_after = Pt(6)
    return tf


def add_placeholder_box(slide, text, left, top, width, height, fill_color=RGBColor(0x1A, 0x1F, 0x2E), border_color=ACCENT_BLUE):
    """스크린샷 들어갈 자리 표시 박스"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(2)
    shape.line.dash_style = 2  # dash
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(14)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER
    return shape


def add_orange_underline(slide, left, top, width):
    """주황색 밑줄 라인"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Pt(4)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_ORANGE
    shape.line.fill.background()


def add_feature_card(slide, title, desc, left, top, width=Inches(3.5), height=Inches(2.2), accent=ACCENT_BLUE):
    """기능 카드"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x16, 0x1B, 0x2E)
    shape.line.color.rgb = accent
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.2)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = accent
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(14)
    p2.font.color.rgb = LIGHT_GRAY
    p2.space_before = Pt(8)


# ════════════════════════════════════════════
# SLIDE 1: 표지
# ════════════════════════════════════════════
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_dark_bg(slide1)

add_title_text(slide1, "BEGIN AGAIN :", Inches(0.8), Inches(0.8), Inches(8), Inches(0.8),
               size=Pt(48), color=WHITE, bold=True)
add_body_text(slide1, "고질적인 현장의 문제를 다시 풀다.", Inches(0.8), Inches(1.6), Inches(8), Inches(0.6),
              size=Pt(24), color=LIGHT_GRAY)

# 서브타이틀
add_multiline_text(slide1, [
    ("Beyond Engineering", False, LIGHT_GRAY, Pt(20)),
    ("GS Intelligent Navigation", False, LIGHT_GRAY, Pt(20)),
    ("Again", False, ACCENT_ORANGE, Pt(20)),
], Inches(0.8), Inches(2.5), Inches(6), Inches(1.5))

# 팀 멤버
members = [
    ("Gon 정승환", "인천종합에너지"),
    ("L 이재현", "GS스포츠"),
    ("Hani 김한희", "GS에너지"),
    ("Yunie 이상윤", "GS엔텍"),
    ("Connor 이성규", "위드인천에너지"),
    ("Eddy 노엘", "GS칼텍스"),
]
for i, (name, company) in enumerate(members):
    col = i % 3
    row = i // 3
    x = Inches(0.8) + col * Inches(2.5)
    y = Inches(4.5) + row * Inches(1.2)
    add_multiline_text(slide1, [
        (name, True, WHITE, Pt(14)),
        (company, False, LIGHT_GRAY, Pt(12)),
    ], x, y, Inches(2.2), Inches(0.8))

# 우측 안내
add_placeholder_box(slide1, "[캐릭터 마스코트 이미지]\nGLANCE 로고 + 스마트폰 목업",
                    Inches(8.5), Inches(1.5), Inches(4), Inches(5),
                    fill_color=RGBColor(0x12, 0x18, 0x2A))


# ════════════════════════════════════════════
# SLIDE 2: P&ID란?
# ════════════════════════════════════════════
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide2)
add_tag(slide2, "P&ID란?", Inches(0.5), Inches(0.4))

add_title_text(slide2, 'P&ID "어떤 설비가 어디에 있고, 배관이 어떻게 연결되며,\n밸브·계측기를 보여주는 지도"',
               Inches(0.5), Inches(1.2), Inches(12), Inches(1.0),
               size=Pt(28), color=DARK_GRAY)
add_orange_underline(slide2, Inches(0.5), Inches(1.15), Inches(12))

add_placeholder_box(slide2, "[현장 설비 사진]\n배관, 밸브가 보이는 실제 발전소 현장",
                    Inches(0.5), Inches(2.5), Inches(5.8), Inches(4.2),
                    fill_color=RGBColor(0xE8, 0xE8, 0xE8), border_color=MEDIUM_GRAY)

add_placeholder_box(slide2, "[P&ID 도면 이미지]\n실제 시공상세도 (배관계장도)",
                    Inches(6.8), Inches(2.5), Inches(5.8), Inches(4.2),
                    fill_color=RGBColor(0xE8, 0xE8, 0xE8), border_color=MEDIUM_GRAY)


# ════════════════════════════════════════════
# SLIDE 3: 비유 (도로지도 vs P&ID)
# ════════════════════════════════════════════
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide3)
add_tag(slide3, "비유하면?!", Inches(0.5), Inches(0.4))

add_title_text(slide3, 'P&ID는 \'발전소의 아주 복잡한 도로지도\'입니다.',
               Inches(0.5), Inches(1.5), Inches(12), Inches(0.8),
               size=Pt(32), color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

add_placeholder_box(slide3, "[전국도로지도 이미지]\n옛날 아빠 차의 필수품: 복잡한 전국도로지도",
                    Inches(0.5), Inches(2.8), Inches(5.8), Inches(3.8),
                    fill_color=RGBColor(0xE8, 0xE8, 0xE8), border_color=MEDIUM_GRAY)

add_placeholder_box(slide3, "[P&ID 도면 이미지]\n발전소의 필수품: 복잡한 P&ID 도면",
                    Inches(6.8), Inches(2.8), Inches(5.8), Inches(3.8),
                    fill_color=RGBColor(0xE8, 0xE8, 0xE8), border_color=MEDIUM_GRAY)


# ════════════════════════════════════════════
# SLIDE 4: 문제정의
# ════════════════════════════════════════════
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide4)
add_tag(slide4, "문제정의", Inches(0.5), Inches(0.4))

add_body_text(slide4, "INTECO 크루 GON의 현업시절 풀고싶었던 문제", Inches(0.5), Inches(1.1), Inches(12), Inches(0.5),
              size=Pt(16), color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

add_title_text(slide4, '작업할 때 마다 찾는 두껍고 방대한 P&ID',
               Inches(0.5), Inches(1.6), Inches(12), Inches(0.8),
               size=Pt(36), color=ACCENT_ORANGE, alignment=PP_ALIGN.CENTER)
add_orange_underline(slide4, Inches(1.5), Inches(1.55), Inches(10))

add_body_text(slide4, '"밸브 및 설비를 찾는게 어렵고 시간이 오래걸려요."',
              Inches(0.5), Inches(2.5), Inches(12), Inches(0.5),
              size=Pt(22), color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

add_placeholder_box(slide4, "[문제정의 만화 5컷]\nDCS 알람 → P&ID 찾기 → 현장 확인 → 도면 vs 실물 괴리",
                    Inches(1), Inches(3.3), Inches(11), Inches(3.5),
                    fill_color=RGBColor(0xE8, 0xE8, 0xE8), border_color=MEDIUM_GRAY)


# ════════════════════════════════════════════
# SLIDE 5: 현업 여정맵
# ════════════════════════════════════════════
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide5)
add_tag(slide5, "현업 여정맵", Inches(0.5), Inches(0.4))

add_title_text(slide5, '기존에 현업은 어떻게 일하고 있었을까요?',
               Inches(0.5), Inches(1.2), Inches(12), Inches(0.8),
               size=Pt(32), color=ACCENT_ORANGE, alignment=PP_ALIGN.CENTER)

add_body_text(slide5, "점검 대상 위치 파악할 때 오래걸림", Inches(0.5), Inches(2.0), Inches(12), Inches(0.5),
              size=Pt(20), color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

# 여정맵 단계
stages = ["현장문제발생\n(감지)", "현장 이동전\n사전 확인", "현장 확인", "현장 조치", "사무실 복귀"]
colors = [RGBColor(0xFF, 0xE0, 0xB2), RGBColor(0xFF, 0x8C, 0x00), RGBColor(0xFF, 0x8C, 0x00),
          RGBColor(0x4C, 0xAF, 0x50), RGBColor(0x4C, 0xAF, 0x50)]
text_colors = [DARK_GRAY, WHITE, WHITE, WHITE, WHITE]

for i, (stage, bg_color, txt_color) in enumerate(zip(stages, colors, text_colors)):
    x = Inches(0.8) + i * Inches(2.4)
    shape = slide5.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(3.0), Inches(2.1), Inches(0.9)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = stage
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = txt_color
    p.alignment = PP_ALIGN.CENTER

# 페인포인트 박스
pain_box = slide5.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.3), Inches(4), Inches(1.8)
)
pain_box.fill.solid()
pain_box.fill.fore_color.rgb = RGBColor(0xFF, 0xEB, 0xEE)
pain_box.line.color.rgb = RGBColor(0xFF, 0x00, 0x00)
pain_box.line.width = Pt(2)
pain_box.line.dash_style = 2
tf = pain_box.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.15)
tf.margin_top = Inches(0.1)
p = tf.paragraphs[0]
p.text = "페인포인트"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = RGBColor(0xFF, 0x00, 0x00)
p2 = tf.add_paragraph()
p2.text = "1. P&ID 검색 시간 소요"
p2.font.size = Pt(14)
p2.font.color.rgb = DARK_GRAY
p3 = tf.add_paragraph()
p3.text = "2. 현장 위치 파악에 시간 소요"
p3.font.size = Pt(14)
p3.font.color.rgb = DARK_GRAY
p4 = tf.add_paragraph()
p4.text = "3. 여러 시스템 왔다갔다 비효율"
p4.font.size = Pt(14)
p4.font.color.rgb = DARK_GRAY


# ════════════════════════════════════════════
# SLIDE 6: 인사이트
# ════════════════════════════════════════════
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide6)
add_tag(slide6, "인사이트", Inches(0.5), Inches(0.4))

add_title_text(slide6, '현업과의 인터뷰 중 우리가 발견한 새로운 인사이트',
               Inches(0.5), Inches(1.2), Inches(12), Inches(0.8),
               size=Pt(30), color=ACCENT_ORANGE, alignment=PP_ALIGN.CENTER)

# 신입 섹션
add_body_text(slide6, "신입 / 저연차 직원", Inches(0.8), Inches(2.3), Inches(5), Inches(0.5),
              size=Pt(22), color=ACCENT_BLUE, bold=True)

junior_quotes = [
    '"신입이라 밸브모양을 몰라서 구글에 밸브 이미지를 검색해봤어요."',
    '"P&ID에서 보면 밸브가 너무 작게 보여서 구분하기 힘들어요."',
    '"일부 설비는 태그가 없어 현장에서 설비 위치를 찾는 데 시간이 걸립니다."',
]
for i, q in enumerate(junior_quotes):
    add_body_text(slide6, q, Inches(0.8), Inches(2.9 + i * 0.55), Inches(11), Inches(0.5),
                  size=Pt(15), color=DARK_GRAY)

# 고연차 섹션
add_body_text(slide6, "고연차 직원", Inches(0.8), Inches(4.8), Inches(5), Inches(0.5),
              size=Pt(22), color=ACCENT_ORANGE, bold=True)

senior_quotes = [
    '"AI로 많은 문제들을 해결하다보니 여러 앱들을 옮겨다니면서 확인을 해야한다."',
    '"인계일지에 내용은 기록이 되어있는데 현장에서 그걸 못보고 조치하다가 사고가 난 적 있었어요."',
]
for i, q in enumerate(senior_quotes):
    add_body_text(slide6, q, Inches(0.8), Inches(5.4 + i * 0.55), Inches(11), Inches(0.5),
                  size=Pt(15), color=DARK_GRAY)


# ════════════════════════════════════════════
# SLIDE 7: 목표 확장 (검색 → 통합 솔루션)
# ════════════════════════════════════════════
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide7)
add_tag(slide7, "인사이트 기반 개선", Inches(0.5), Inches(0.4))

add_title_text(slide7, '프로젝트의 목표를 "검색을 잘되게 하자" 에서',
               Inches(0.5), Inches(1.3), Inches(12), Inches(0.8),
               size=Pt(28), color=ACCENT_ORANGE, alignment=PP_ALIGN.CENTER)
add_orange_underline(slide7, Inches(1), Inches(1.25), Inches(11))

add_title_text(slide7, '"한 화면 통합 + 실시간 시스템 연동" 으로 확장',
               Inches(0.5), Inches(2.0), Inches(12), Inches(0.8),
               size=Pt(28), color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

# Before 박스
before_box = slide7.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(3.2), Inches(4.5), Inches(3.5)
)
before_box.fill.solid()
before_box.fill.fore_color.rgb = RGBColor(0xE8, 0xE8, 0xE8)
before_box.line.fill.background()
tf = before_box.text_frame
tf.word_wrap = True
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.text = "Search Tool"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = DARK_GRAY
p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = "\nP&ID를 아는 사람\n기준의 OCR/검색 툴."
p2.font.size = Pt(16)
p2.font.color.rgb = MEDIUM_GRAY
p2.alignment = PP_ALIGN.CENTER

# 화살표 텍스트
add_title_text(slide7, "GLANCE\n→", Inches(5.5), Inches(4.0), Inches(2), Inches(1.5),
               size=Pt(32), color=ACCENT_ORANGE, alignment=PP_ALIGN.CENTER)

# After 박스
after_box = slide7.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.8), Inches(3.2), Inches(4.8), Inches(3.5)
)
after_box.fill.solid()
after_box.fill.fore_color.rgb = ACCENT_BLUE
after_box.line.fill.background()
tf = after_box.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.3)
tf.margin_top = Inches(0.3)
p = tf.paragraphs[0]
p.text = "Integrated Solution"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = WHITE

features = [
    "신입 기준으로 설비를 찾고 (P&ID)",
    "모양을 이해하고 (3D 모델 + 사진)",
    "현장을 미리 보고 (360° 로드뷰)",
    "과거의 경험을 함께 보며 (정비이력)",
    "한 화면에서 판단하는 통합 솔루션.",
]
for feat in features:
    p2 = tf.add_paragraph()
    p2.text = feat
    p2.font.size = Pt(14)
    p2.font.color.rgb = WHITE
    p2.space_before = Pt(6)


# ════════════════════════════════════════════
# SLIDE 8: GLANCE 소개 - 핵심 기능 (NEW)
# ════════════════════════════════════════════
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide8)
add_tag(slide8, "GLANCE 소개", Inches(0.5), Inches(0.4), color=ACCENT_BLUE)

add_title_text(slide8, '솔루션 : GLANCE / 뜻 : 훑어보다!',
               Inches(0.5), Inches(1.2), Inches(12), Inches(0.8),
               size=Pt(36), color=ACCENT_ORANGE, alignment=PP_ALIGN.CENTER)
add_orange_underline(slide8, Inches(1), Inches(1.15), Inches(11))

add_body_text(slide8, "밸브 정보를 한눈에(Glance) 찾아주는 스마트 통합 솔루션",
              Inches(0.5), Inches(2.0), Inches(12), Inches(0.5),
              size=Pt(20), color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# 6개 핵심 기능 카드
features_data = [
    ("GenAI Search", "자연어 검색으로 위치 점프\nKKS 코드 또는 자연어로\n7,542개 컴포넌트 즉시 검색", ACCENT_BLUE),
    ("3D Model Preview", "밸브 종류별 3D 모델 표시\n회전 가능한 인터랙티브 뷰\n신입도 밸브 형상 즉시 파악", ACCENT_BLUE),
    ("360° Roadview", "파노라마 현장 사진 연동\n226개 씬으로 현장 가상 투어\n도면 ↔ 현장 즉시 연결", ACCENT_GREEN),
    ("EAM API 연동", "HxGN EAM 정비이력 실시간 조회\n작업오더, 승인상태, 담당자\n실제 데이터 — 목업 아님!", ACCENT_ORANGE),
    ("인계일지 통합", "근무조별 인계사항 즉시 확인\n밸브 조회 시 관련 이력 자동 표시\n암묵지의 디지털화", ACCENT_ORANGE),
    ("밸브/계기 사전", "P&ID 심볼 ↔ 실물 매칭\n기술사양 (압력, 온도, 크기)\n제조사 ID까지 원클릭 확인", ACCENT_BLUE),
]

for i, (title, desc, color) in enumerate(features_data):
    col = i % 3
    row = i // 3
    x = Inches(0.5) + col * Inches(4.2)
    y = Inches(2.9) + row * Inches(2.3)
    add_feature_card(slide8, title, desc, x, y, Inches(3.9), Inches(2.0), accent=color)


# ════════════════════════════════════════════
# SLIDE 9: GLANCE 메인 화면 스크린샷
# ════════════════════════════════════════════
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide9)
add_tag(slide9, "개선 내용", Inches(0.5), Inches(0.4), color=ACCENT_BLUE)

add_title_text(slide9, 'GLANCE 메인 화면',
               Inches(0.5), Inches(1.2), Inches(12), Inches(0.8),
               size=Pt(32), color=WHITE, alignment=PP_ALIGN.CENTER)

add_placeholder_box(slide9, "[GLANCE 메인 화면 스크린샷]\nP&ID 도면 뷰어 + 검색바 + 왼쪽 정보카드 3개\n(밸브/계기 사전, 기술사양, 정비이력 매칭)\n+ 하단 컴포넌트 수 카드\n\nlocalhost:3010 에서 캡처",
                    Inches(1), Inches(2.2), Inches(11), Inches(4.8))


# ════════════════════════════════════════════
# SLIDE 10: 밸브 상세 패널
# ════════════════════════════════════════════
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide10)
add_tag(slide10, "개선 내용", Inches(0.5), Inches(0.4), color=ACCENT_BLUE)

add_title_text(slide10, '밸브를 보는 순간',
               Inches(0.5), Inches(1.0), Inches(12), Inches(0.6),
               size=Pt(36), color=ACCENT_ORANGE, alignment=PP_ALIGN.CENTER)
add_orange_underline(slide10, Inches(3), Inches(0.95), Inches(7))

add_body_text(slide10, "3D 모델 + 기술사양 + 정비이력 + 인계일지 + 로드뷰를 한 화면에서",
              Inches(0.5), Inches(1.7), Inches(12), Inches(0.5),
              size=Pt(18), color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_placeholder_box(slide10, "[밸브 상세패널 스크린샷]\n좌: P&ID 도면 + 밸브 위치 마커\n우: 상세패널 (3D 모델, 기술사양, 정비이력, 인계일지)\n\nlocalhost:3010 에서 밸브 클릭 후 캡처",
                    Inches(0.5), Inches(2.5), Inches(7.5), Inches(4.5))

# 우측 키포인트
keypoints = [
    ("One-Screen Intelligence", "앱 이동 없이\n한 화면에서 판단 종료", ACCENT_BLUE),
    ("Digital Tacit Knowledge", "불안감이\n확신으로 바뀝니다", ACCENT_GREEN),
    ("Real-time EAM Data", "HxGN EAM API 실시간\n정비이력 연동 — 목업 아님", ACCENT_ORANGE),
]
for i, (title, desc, color) in enumerate(keypoints):
    y = Inches(2.6) + i * Inches(1.5)
    add_body_text(slide10, title, Inches(8.5), y, Inches(4), Inches(0.4),
                  size=Pt(18), color=color, bold=True)
    add_body_text(slide10, desc, Inches(8.5), y + Inches(0.35), Inches(4), Inches(0.8),
                  size=Pt(14), color=LIGHT_GRAY)


# ════════════════════════════════════════════
# SLIDE 11: 360° 로드뷰 (NEW FEATURE!)
# ════════════════════════════════════════════
slide11 = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide11)
add_tag(slide11, "NEW 로드뷰", Inches(0.5), Inches(0.4), color=ACCENT_GREEN)

add_title_text(slide11, '도면에서 → 현장으로 즉시 연결',
               Inches(0.5), Inches(1.0), Inches(12), Inches(0.6),
               size=Pt(36), color=ACCENT_GREEN, alignment=PP_ALIGN.CENTER)
add_body_text(slide11, "360° 파노라마 로드뷰로 현장을 가상 투어",
              Inches(0.5), Inches(1.7), Inches(12), Inches(0.5),
              size=Pt(20), color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_placeholder_box(slide11, "[로드뷰 화면 스크린샷]\n360° 파노라마 뷰어\n핫스팟 네비게이션 화살표\n\nlocalhost:3010/roadview 에서 캡처",
                    Inches(0.5), Inches(2.5), Inches(7.5), Inches(4.5))

# 우측 스펙
specs = [
    "226개 파노라마 씬",
    "ST동 + 발전소 전체 투어",
    "핫스팟으로 씬 간 이동",
    "밸브 상세패널에서 바로 연결",
    "수평 보정 자동 적용",
    "도면 ↔ 현장 괴리 해소",
]
for i, spec in enumerate(specs):
    y = Inches(2.8) + i * Inches(0.6)
    add_body_text(slide11, f"▸ {spec}", Inches(8.5), y, Inches(4), Inches(0.5),
                  size=Pt(16), color=WHITE)


# ════════════════════════════════════════════
# SLIDE 12: 시스템 아키텍처 (NEW)
# ════════════════════════════════════════════
slide12 = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide12)
add_tag(slide12, "시스템 구조", Inches(0.5), Inches(0.4), color=ACCENT_BLUE)

add_title_text(slide12, 'GLANCE 시스템 아키텍처',
               Inches(0.5), Inches(1.2), Inches(12), Inches(0.8),
               size=Pt(32), color=WHITE, alignment=PP_ALIGN.CENTER)

# 프론트엔드 박스
fe_box = slide12.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(2.5), Inches(4), Inches(4.2)
)
fe_box.fill.solid()
fe_box.fill.fore_color.rgb = RGBColor(0x16, 0x1B, 0x2E)
fe_box.line.color.rgb = ACCENT_BLUE
fe_box.line.width = Pt(2)
tf = fe_box.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.2)
tf.margin_top = Inches(0.15)
p = tf.paragraphs[0]
p.text = "Frontend"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = ACCENT_BLUE
items = ["Next.js 15 (App Router)", "React Three Fiber (3D)", "Pannellum (360° 뷰어)",
         "Tailwind CSS", "GenAI 자연어 검색"]
for item in items:
    p2 = tf.add_paragraph()
    p2.text = f"  ▸ {item}"
    p2.font.size = Pt(14)
    p2.font.color.rgb = LIGHT_GRAY
    p2.space_before = Pt(6)

# 중앙 화살표
add_title_text(slide12, "←→", Inches(4.8), Inches(4.0), Inches(1.5), Inches(0.8),
               size=Pt(36), color=ACCENT_ORANGE, alignment=PP_ALIGN.CENTER)

# API 연동 박스
api_box = slide12.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.5), Inches(2.5), Inches(3), Inches(4.2)
)
api_box.fill.solid()
api_box.fill.fore_color.rgb = RGBColor(0x16, 0x1B, 0x2E)
api_box.line.color.rgb = ACCENT_ORANGE
api_box.line.width = Pt(2)
tf = api_box.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.2)
tf.margin_top = Inches(0.15)
p = tf.paragraphs[0]
p.text = "API 연동"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = ACCENT_ORANGE
api_items = ["HxGN EAM\n(정비이력 실시간)", "REST API\n(OAuth2 인증)", "JSON 데이터\n(밸브사양/도면)"]
for item in api_items:
    p2 = tf.add_paragraph()
    p2.text = f"  ▸ {item}"
    p2.font.size = Pt(13)
    p2.font.color.rgb = LIGHT_GRAY
    p2.space_before = Pt(8)

# 데이터 박스
data_box = slide12.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10), Inches(2.5), Inches(2.8), Inches(4.2)
)
data_box.fill.solid()
data_box.fill.fore_color.rgb = RGBColor(0x16, 0x1B, 0x2E)
data_box.line.color.rgb = ACCENT_GREEN
data_box.line.width = Pt(2)
tf = data_box.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.15)
tf.margin_top = Inches(0.15)
p = tf.paragraphs[0]
p.text = "데이터"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = ACCENT_GREEN
data_items = ["P&ID 도면 이미지", "7,542개 컴포넌트", "226개 파노라마 씬", "밸브 기술사양", "인계일지"]
for item in data_items:
    p2 = tf.add_paragraph()
    p2.text = f"  ▸ {item}"
    p2.font.size = Pt(13)
    p2.font.color.rgb = LIGHT_GRAY
    p2.space_before = Pt(6)


# ════════════════════════════════════════════
# SLIDE 13: 시연영상
# ════════════════════════════════════════════
slide13 = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide13)
add_tag(slide13, "시연영상", Inches(0.5), Inches(0.4), color=ACCENT_ORANGE)

add_placeholder_box(slide13, "[시연 영상 또는 GIF]\n\n1. 검색창에 밸브 태그 입력\n2. P&ID에서 위치 자동 점프\n3. 상세패널 (3D + 기술사양 + 정비이력)\n4. 로드뷰 버튼 → 360° 현장 확인\n5. 인계일지 확인\n\n최신 버전으로 재촬영 필요!",
                    Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.2))


# ════════════════════════════════════════════
# SLIDE 14: Before → After 결과
# ════════════════════════════════════════════
slide14 = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide14)
add_tag(slide14, "결과", Inches(0.5), Inches(0.4))

add_title_text(slide14, '5초 만에 도면 및 위치 확인하고, 이력까지 확인가능',
               Inches(0.5), Inches(1.2), Inches(12), Inches(0.8),
               size=Pt(30), color=ACCENT_ORANGE, alignment=PP_ALIGN.CENTER)
add_orange_underline(slide14, Inches(1), Inches(1.15), Inches(11))

# Before/After 비교
before = slide14.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(2.2), Inches(5.8), Inches(4.5)
)
before.fill.solid()
before.fill.fore_color.rgb = RGBColor(0xFF, 0xEB, 0xEE)
before.line.fill.background()
tf = before.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.3)
tf.margin_top = Inches(0.3)
p = tf.paragraphs[0]
p.text = "Before"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = RGBColor(0xEF, 0x44, 0x44)

before_items = [
    "두꺼운 P&ID 책 뒤적거리기",
    "밸브 모양 몰라서 구글 검색",
    "여러 시스템 왔다갔다",
    "인계사항 확인 누락 → 사고 위험",
    "현장 위치 파악에 수십분 소요",
]
for item in before_items:
    p2 = tf.add_paragraph()
    p2.text = f"✗  {item}"
    p2.font.size = Pt(15)
    p2.font.color.rgb = DARK_GRAY
    p2.space_before = Pt(8)

after = slide14.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(2.2), Inches(5.8), Inches(4.5)
)
after.fill.solid()
after.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
after.line.fill.background()
tf = after.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.3)
tf.margin_top = Inches(0.3)
p = tf.paragraphs[0]
p.text = "After (GLANCE)"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = RGBColor(0x16, 0xA3, 0x4A)

after_items = [
    "3초 검색으로 P&ID 위치 즉시 확인",
    "3D 모델로 밸브 형상 바로 파악",
    "360° 로드뷰로 현장 미리 확인",
    "정비이력 + 인계일지 자동 표시",
    "한 화면에서 판단 완료",
]
for item in after_items:
    p2 = tf.add_paragraph()
    p2.text = f"✓  {item}"
    p2.font.size = Pt(15)
    p2.font.color.rgb = DARK_GRAY
    p2.space_before = Pt(8)


# ════════════════════════════════════════════
# SLIDE 15: 정량 성과 (NEW)
# ════════════════════════════════════════════
slide15 = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide15)
add_tag(slide15, "정량 성과", Inches(0.5), Inches(0.4), color=ACCENT_ORANGE)

add_title_text(slide15, 'GLANCE 고도화 현황',
               Inches(0.5), Inches(1.2), Inches(12), Inches(0.8),
               size=Pt(36), color=WHITE, alignment=PP_ALIGN.CENTER)

# 숫자 카드들
metrics = [
    ("7,542", "검색 가능 컴포넌트", ACCENT_BLUE),
    ("226", "파노라마 씬 (360°)", ACCENT_GREEN),
    ("6종", "3D 밸브 모델", ACCENT_BLUE),
    ("실시간", "HxGN EAM API 연동", ACCENT_ORANGE),
]

for i, (number, label, color) in enumerate(metrics):
    x = Inches(0.5) + i * Inches(3.2)
    # 숫자
    box = slide15.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.8), Inches(2.9), Inches(2.5)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0x16, 0x1B, 0x2E)
    box.line.color.rgb = color
    box.line.width = Pt(2)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = number
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.size = Pt(16)
    p2.font.color.rgb = LIGHT_GRAY
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(10)

# Phase 진행 상황
add_body_text(slide15, "Phase 1 완료: P&ID 검색 + 3D 모델 + 밸브 사전 + 기술사양",
              Inches(0.5), Inches(5.8), Inches(12), Inches(0.4),
              size=Pt(16), color=ACCENT_BLUE)
add_body_text(slide15, "Phase 2 완료: HxGN EAM 연동 + 360° 로드뷰 + 인계일지 통합",
              Inches(0.5), Inches(6.2), Inches(12), Inches(0.4),
              size=Pt(16), color=ACCENT_GREEN)


# ════════════════════════════════════════════
# SLIDE 16: 향후 계획
# ════════════════════════════════════════════
slide16 = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide16)
add_tag(slide16, "향후 계획", Inches(0.5), Inches(0.4))

add_title_text(slide16, 'Next Step!',
               Inches(0.5), Inches(1.2), Inches(12), Inches(0.8),
               size=Pt(40), color=ACCENT_ORANGE, alignment=PP_ALIGN.CENTER)

# Phase 표
phases = [
    ("Phase 1 ✓", "P&ID 검색 + 3D 밸브 모델\n밸브/계기 사전 + 기술사양", "완료", ACCENT_BLUE),
    ("Phase 2 ✓", "HxGN EAM 정비이력 API 연동\n360° 로드뷰 + 인계일지 통합", "완료", ACCENT_GREEN),
    ("Phase 3", "PI System 실시간 데이터 연동\n사내 NAS 파일서버 이미지 관리\nAI 기반 이상 감지 알림", "예정", ACCENT_ORANGE),
]

for i, (phase, desc, status, color) in enumerate(phases):
    x = Inches(0.5) + i * Inches(4.2)
    box = slide16.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.5), Inches(3.9), Inches(3.8)
    )
    box.fill.solid()
    if status == "완료":
        box.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9) if color == ACCENT_GREEN else RGBColor(0xE3, 0xF2, 0xFD)
    else:
        box.fill.fore_color.rgb = RGBColor(0xFF, 0xF3, 0xE0)
    box.line.color.rgb = color
    box.line.width = Pt(2)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.2)
    p = tf.paragraphs[0]
    p.text = phase
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = color
    p2 = tf.add_paragraph()
    p2.text = ""
    p2.font.size = Pt(8)
    for line in desc.split("\n"):
        p3 = tf.add_paragraph()
        p3.text = f"  ▸ {line}"
        p3.font.size = Pt(14)
        p3.font.color.rgb = DARK_GRAY
        p3.space_before = Pt(6)

    # 상태 배지
    p_status = tf.add_paragraph()
    p_status.text = f"\n[ {status} ]"
    p_status.font.size = Pt(16)
    p_status.font.bold = True
    p_status.font.color.rgb = color
    p_status.alignment = PP_ALIGN.CENTER


# ════════════════════════════════════════════
# SLIDE 17: L&L (팀 소감)
# ════════════════════════════════════════════
slide17 = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide17)

add_title_text(slide17, 'Crucible Hackathon L&L',
               Inches(0.5), Inches(0.5), Inches(12), Inches(0.8),
               size=Pt(36), color=DARK_GRAY)
add_body_text(slide17, "(우리는 무엇을 배웠나요?!)", Inches(0.5), Inches(1.2), Inches(12), Inches(0.5),
              size=Pt(20), color=MEDIUM_GRAY)

quotes = [
    ("현장의 목소리에 답이 있다는 것을\n또 한번 느꼈습니다.", ACCENT_ORANGE),
    ("새로 만드는 것보다,\n흩어진 것을 연결하는 것이 먼저.", ACCENT_BLUE),
    ("NO FEAR — 완벽하지 않아도\n일단 해보는 용기.", ACCENT_GREEN),
    ("익숙해져서 지나쳤던 문제들이\n생각보다 많다.", ACCENT_ORANGE),
]

for i, (quote, color) in enumerate(quotes):
    col = i % 2
    row = i // 2
    x = Inches(0.5) + col * Inches(6.3)
    y = Inches(2.2) + row * Inches(2.3)
    box = slide17.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.8), Inches(1.8)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = color
    box.line.width = Pt(2)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = f'"{quote}"'
    p.font.size = Pt(18)
    p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.LEFT

add_placeholder_box(slide17, "[팀 사진]",
                    Inches(4.5), Inches(5.5), Inches(4), Inches(1.5),
                    fill_color=RGBColor(0xE8, 0xE8, 0xE8), border_color=MEDIUM_GRAY)


# ════════════════════════════════════════════
# SLIDE 18: Q&A
# ════════════════════════════════════════════
slide18 = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide18)

add_title_text(slide18, 'Q & A',
               Inches(0.5), Inches(2.5), Inches(12), Inches(1.5),
               size=Pt(72), color=WHITE, alignment=PP_ALIGN.CENTER)

add_body_text(slide18, "GLANCE — 밸브 정보를 한눈에 찾아주는 스마트 통합 솔루션",
              Inches(0.5), Inches(4.5), Inches(12), Inches(0.5),
              size=Pt(20), color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_body_text(slide18, "감사합니다",
              Inches(0.5), Inches(5.5), Inches(12), Inches(0.5),
              size=Pt(24), color=ACCENT_ORANGE, alignment=PP_ALIGN.CENTER)


# ── 저장 ──
output_path = os.path.join(os.path.dirname(__file__), "..", "GLANCE_발표자료_v2.pptx")
prs.save(output_path)
print(f"발표자료 저장 완료: {os.path.abspath(output_path)}")
print(f"총 {len(prs.slides)} 슬라이드")
