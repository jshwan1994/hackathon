# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color palette
BG_DARK = RGBColor(0x0F, 0x17, 0x2A)
BG_CARD = RGBColor(0x1A, 0x25, 0x3C)
ACCENT = RGBColor(0x38, 0xBD, 0xF8)
ACCENT2 = RGBColor(0x22, 0xD3, 0xEE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x94, 0xA3, 0xB8)
LIGHT = RGBColor(0xCB, 0xD5, 0xE1)
ORANGE = RGBColor(0xFB, 0x92, 0x3C)
GREEN = RGBColor(0x4A, 0xDE, 0x80)
RED = RGBColor(0xF8, 0x71, 0x71)
YELLOW = RGBColor(0xFA, 0xCC, 0x15)
BG_HIGHLIGHT = RGBColor(0x15, 0x2B, 0x44)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    shape.adjustments[0] = 0.05
    return shape


def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "\ub9d1\uc740 \uace0\ub515"
    p.alignment = align
    return txBox


def add_multiline(slide, left, top, width, height, lines, size=16, color=LIGHT, line_spacing=1.3):
    """Add multi-line text with individual line control"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, txt_color, txt_bold, txt_size) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(txt_size if txt_size else size)
        run.font.color.rgb = txt_color if txt_color else color
        run.font.bold = txt_bold
        run.font.name = "\ub9d1\uc740 \uace0\ub515"
    return txBox


def add_divider(slide, y=Inches(1.05)):
    shape = add_rect(slide, Inches(0.8), y, Inches(11.7), Pt(1), RGBColor(0x33, 0x44, 0x55))
    return shape


def add_section_header(slide, number, title):
    add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
             f"{number}  {title}", size=30, color=ACCENT, bold=True)
    add_divider(slide)


# ============================================================
# SLIDE 1: Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_rect(slide, Inches(0), Inches(3.1), Inches(13.333), Pt(3), ACCENT)
add_text(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1),
         "GLANCE", size=54, color=ACCENT, bold=True)
add_text(slide, Inches(1.5), Inches(2.45), Inches(10), Inches(0.7),
         "AI Code Generation \uae30\ubc18 \ud604\uc7a5 \ubc38\ube0c \uad00\ub9ac \ud504\ub85c\ub355\ud2b8", size=28, color=WHITE, bold=True)
add_text(slide, Inches(1.5), Inches(3.6), Inches(10), Inches(0.5),
         "52G \ud611\uc758\uccb4 \u2014 AI Native Product \uac70\ubc84\ub10c\uc2a4 \uc218\ub9bd\uc744 \uc704\ud55c \uc0ac\ub840 \ubc1c\ud45c", size=18, color=GRAY)
add_text(slide, Inches(1.5), Inches(4.2), Inches(10), Inches(0.5),
         "2026. 03. 12", size=16, color=GRAY)


# ============================================================
# SLIDE 2: Product Overview
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_section_header(slide, "01", "\ud504\ub85c\ub355\ud2b8 \uac1c\uc694")

add_text(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.8),
         "\ubc1c\uc804\uc18c \ud604\uc7a5\uc758 \ubc38\ube0c\xb7\uacc4\uae30\ub97c P&ID \ub3c4\uba74 \uc704\uc5d0\uc11c \uac80\uc0c9\ud558\uace0,\n360\xb0 \ud30c\ub178\ub77c\ub9c8 \ub85c\ub4dc\ubdf0\ub85c \ud604\uc7a5\uc744 \ud655\uc778\ud558\uba70, EAM \uc815\ube44\uc774\ub825\uae4c\uc9c0 \ud55c \ud654\uba74\uc5d0\uc11c \uc870\ud68c\ud558\ub294 \uc6f9 \uc560\ud50c\ub9ac\ucf00\uc774\uc158",
         size=18, color=LIGHT)

features = [
    ("P&ID \ub3c4\uba74 \ubdf0\uc5b4", "\ub3c4\uba74 \uc704 \ubc38\ube0c \uc704\uce58\n\uac80\uc0c9 \ubc0f \ub9c8\ucee4 \ud45c\uc2dc"),
    ("360\xb0 \ub85c\ub4dc\ubdf0", "226\uac1c \uc528 \ud30c\ub178\ub77c\ub9c8\n\ud604\uc7a5 \uc2e4\uc0ac \ud0d0\uc0c9"),
    ("EAM \uc815\ube44\uc774\ub825", "HxGN EAM \uc5f0\ub3d9\n\uc2e4\uc2dc\uac04 \uc774\ub825 \uc870\ud68c"),
    ("\ubc38\ube0c/\uacc4\uae30 \uc0ac\uc804", "\uce74\ud14c\uace0\ub9ac\ubcc4 \ubd84\ub958\n\uaca9\ub9ac \uc808\ucc28 \uc548\ub0b4"),
    ("3D \ubc38\ube0c \ubaa8\ub378", "Three.js \uae30\ubc18\n\uc785\uccb4 \uc2dc\uac01\ud654"),
]

card_w = Inches(2.1)
card_h = Inches(2.2)
start_x = Inches(0.8)
gap = Inches(0.25)

for i, (title, desc) in enumerate(features):
    x = start_x + i * (card_w + gap)
    y = Inches(2.5)
    add_shape_bg(slide, x, y, card_w, card_h, BG_CARD)
    add_text(slide, x + Inches(0.2), y + Inches(0.3), card_w - Inches(0.4), Inches(0.5),
             title, size=18, color=ACCENT2, bold=True)
    add_text(slide, x + Inches(0.2), y + Inches(0.9), card_w - Inches(0.4), Inches(1),
             desc, size=14, color=LIGHT)

stats = [("\uac1c\ubc1c \uae30\uac04", "\uc57d 2\uac1c\uc6d4"), ("\ucee4\ubc0b \uc218", "48\ud68c"), ("\uc2dc\uc791\uc77c", "2026.01.21"), ("\uac1c\ubc1c \ub3c4\uad6c", "Claude Code")]
stat_w = Inches(2.6)
for i, (label, val) in enumerate(stats):
    x = Inches(0.8) + i * (stat_w + Inches(0.2))
    y = Inches(5.2)
    add_text(slide, x, y, stat_w, Inches(0.35), label, size=13, color=GRAY)
    add_text(slide, x, y + Inches(0.35), stat_w, Inches(0.4), val, size=20, color=WHITE, bold=True)


# ============================================================
# SLIDE 3: Architecture
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_section_header(slide, "02", "\uc544\ud0a4\ud14d\ucc98")

tiers = [
    ("Frontend", "Next.js 15 \xb7 React 19 \xb7 Tailwind CSS \xb7 Three.js", Inches(1.4), ACCENT, [
        "P&ID \ub3c4\uba74 \ubdf0\uc5b4 (Canvas)", "360\xb0 \ud30c\ub178\ub77c\ub9c8 \ub85c\ub4dc\ubdf0", "3D \ubc38\ube0c \ubaa8\ub378 \ubdf0\uc5b4", "\ubc38\ube0c \uc0c1\uc138 \ud328\ub110 / \uc0ac\uc804"
    ]),
    ("BFF (API Routes)", "Next.js Server-side API Routes", Inches(3.3), ACCENT2, [
        "/api/maintenance \u2014 EAM \ud504\ub85d\uc2dc", "/api/diaries \u2014 \uc778\uacc4\uc77c\uc9c0 \ud504\ub85d\uc2dc", "/api/roadview-settings \u2014 \uc124\uc815 \uad00\ub9ac"
    ]),
    ("External / Legacy", "\uc0ac\ub0b4 \uc2dc\uc2a4\ud15c \ubc0f \ud074\ub77c\uc6b0\ub4dc \uc11c\ube44\uc2a4", Inches(5.0), ORANGE, [
        "HxGN EAM REST API (\uc815\ube44\uc774\ub825)", "AWS Lambda (\uc778\uacc4\uc77c\uc9c0)", "JSON \ud30c\uc77c (\ubc38\ube0c \ub9c8\uc2a4\ud130 \ub370\uc774\ud130)"
    ]),
]

for title, subtitle, y, color, items in tiers:
    add_shape_bg(slide, Inches(0.8), y, Inches(11.7), Inches(1.5), BG_CARD)
    add_rect(slide, Inches(0.8), y, Pt(5), Inches(1.5), color)
    add_text(slide, Inches(1.2), y + Inches(0.1), Inches(3), Inches(0.4),
             title, size=20, color=color, bold=True)
    add_text(slide, Inches(1.2), y + Inches(0.5), Inches(3.5), Inches(0.3),
             subtitle, size=12, color=GRAY)
    for j, item in enumerate(items):
        col = j // 2
        row = j % 2
        add_text(slide, Inches(5.5) + col * Inches(3.5), y + Inches(0.2) + row * Inches(0.55), Inches(3.5), Inches(0.5),
                 f"\u25b8  {item}", size=13, color=LIGHT)


# ============================================================
# SLIDE 4: AI Code Generation
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_section_header(slide, "03", "AI Code Generation \ud65c\uc6a9")

add_text(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.6),
         "\ucc98\uc74c\ubd80\ud130 \ub05d\uae4c\uc9c0 Claude Code(AI \ucf54\ub4dc \uc0dd\uc131)\ub85c \uac1c\ubc1c \u2014 \uc804\ubb38 \uac1c\ubc1c\uc790 \uc5c6\uc774 \ud604\uc5c5 \ub2f4\ub2f9\uc790\uac00 \uc9c1\uc811 \uad6c\ud604",
         size=18, color=YELLOW, bold=True)

ai_items = [
    ("\ud504\ub85c\uc81d\ud2b8 \uad6c\uc870 \uc124\uacc4", "Next.js \uc2a4\uce90\ud3f4\ub529,\n\ucef4\ud3ec\ub10c\ud2b8 \uc124\uacc4\ub97c\nAI\uc640 \ub300\ud654\ud558\uba70 \uad6c\uc131", "\ucd08\uae30 \uad6c\uc870"),
    ("\ud575\uc2ec \uae30\ub2a5 \uad6c\ud604", "P&ID \ubdf0\uc5b4, 360\xb0 \ud30c\ub178\ub77c\ub9c8,\n3D \ubaa8\ub378 \u2014 \uc804\uccb4 \ucf54\ub4dc\nAI\uac00 \uc0dd\uc131", "\ucf54\uc5b4 \uac1c\ubc1c"),
    ("\ub808\uac70\uc2dc API \uc5f0\ub3d9", "EAM REST API \ud30c\uc2f1,\nAWS Lambda \ud638\ucd9c \ucf54\ub4dc\ub97c\n\ud504\ub86c\ud504\ud2b8\ub85c \uc791\uc131", "API \uc5f0\ub3d9"),
    ("UI/UX \ub514\uc790\uc778", "Tailwind \ubc18\uc751\ud615 UI,\n\uce74\ub4dc \ud6a8\uacfc, \ud14c\ub9c8\ub97c\nAI\uac00 \ub514\uc790\uc778+\ucf54\ub4dc \ub3d9\uc2dc \uc0dd\uc131", "\ub514\uc790\uc778"),
    ("\ub514\ubc84\uae45\xb7\ub9ac\ud329\ud1a0\ub9c1", "\ud30c\ub178\ub77c\ub9c8 \ub4a4\ud2c0\ub9bc \ubc84\uadf8 \ub4f1\n\uc99d\uc0c1 \uc124\uba85 \u2192 AI\uac00\n\uc218\uc815 \ucf54\ub4dc \uc81c\uacf5", "\uc720\uc9c0\ubcf4\uc218"),
]

for i, (title, desc, tag) in enumerate(ai_items):
    x = Inches(0.8) + i * Inches(2.45)
    y = Inches(2.2)
    add_shape_bg(slide, x, y, Inches(2.25), Inches(2.8), BG_CARD)
    add_text(slide, x + Inches(0.15), y + Inches(0.15), Inches(1.8), Inches(0.3),
             tag, size=11, color=ACCENT, bold=True)
    add_text(slide, x + Inches(0.15), y + Inches(0.55), Inches(1.95), Inches(0.5),
             title, size=16, color=WHITE, bold=True)
    add_text(slide, x + Inches(0.15), y + Inches(1.15), Inches(1.95), Inches(1.4),
             desc, size=13, color=LIGHT)

box = add_shape_bg(slide, Inches(0.8), Inches(5.4), Inches(11.7), Inches(1.2), BG_HIGHLIGHT)
add_text(slide, Inches(1.2), Inches(5.55), Inches(11), Inches(0.4),
         "\ud575\uc2ec \ud3ec\uc778\ud2b8", size=14, color=ACCENT, bold=True)
add_text(slide, Inches(1.2), Inches(5.95), Inches(11), Inches(0.5),
         "AI Code Generation \ub355\ubd84\uc5d0 \uc544\uc774\ub514\uc5b4-\ud504\ub85c\ud1a0\ud0c0\uc785-\ud504\ub85c\ub355\ud2b8 \uc0ac\uc774\ud074\uc774 \uadf9\ub2e8\uc801\uc73c\ub85c \uc9e7\uc544\uc9d0.\n\uc774\ub7f0 \uc2dc\ub3c4\ub97c \uc870\uc9c1 \uc548\uc5d0\uc11c \ub354 \ud65c\ubc1c\ud558\uac8c \ud65c\uc6a9\ud558\ub824\uba74 \uc5b4\ub5a4 \uccb4\uacc4\uac00 \ud544\uc694\ud560\uae4c?",
         size=15, color=LIGHT)


# ============================================================
# SLIDE 5: Legacy 연결 경험
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_section_header(slide, "04", "\ub808\uac70\uc2dc \uc5f0\uacb0 \uacbd\ud5d8")

# Left card
add_shape_bg(slide, Inches(0.8), Inches(1.4), Inches(5.6), Inches(5.2), BG_CARD)
add_text(slide, Inches(1.1), Inches(1.6), Inches(5), Inches(0.4),
         "\uac1c\ubc1c \uacfc\uc815\uc5d0\uc11c \ub290\ub080 \uc138 \uac00\uc9c0", size=20, color=ACCENT2, bold=True)

feelings = [
    ("AI\uc758 \uac00\ub2a5\uc131\uc740 \uae30\ub300 \uc774\uc0c1",
     "360\xb0 \ud30c\ub178\ub77c\ub9c8 \ubdf0\uc5b4\ub97c \ub9cc\ub4e4\uc5b4\uc918\ub77c\uace0 \ud588\ub354\ub2c8\nThree.js\ub85c 226\uac1c \uc528 \uc804\ud658\ud558\ub294 \ucf54\ub4dc\ub97c \uc0dd\uc131.\n\uc0c1\uc0c1\ub9cc\uc73c\ub85c \uc2e4\uc81c \ub3d9\uc791\ud558\ub294 \ud504\ub85c\ub355\ud2b8 \uc644\uc131"),
    ("\ub808\uac70\uc2dc \uc5f0\ub3d9 \uc2dc \uc2dc\ud589\ucc29\uc624",
     "EAM API \uc778\uc99d \ubc29\uc2dd, \uad8c\ud55c \uc808\ucc28, \ub370\uc774\ud130 \ud3ec\ub9f7 \ub4f1\n\uc0ac\uc804 \uac00\uc774\ub4dc\uac00 \uc5c6\uc5b4 \ud558\ub098\ud558\ub098 \uc9c1\uc811\n\ubd80\ub52a\ud788\uba70 \uc54c\uc544\ub0c4"),
    ("\uadf8 \ub2e4\uc74c \ub2e8\uacc4\uc758 \uace0\ubbfc",
     "\uae30\ub2a5\uc740 \ub9cc\ub4e4\uc5c8\ub294\ub370, \uc5b4\ub514\uc5d0 \uc5b4\ub5bb\uac8c \ubc30\ud3ec\ud574\uc57c \ud558\ub294\uc9c0\n\uae30\uc900\uc774 \ubcf4\uc774\uc9c0 \uc54a\uc558\ub358 \uacbd\ud5d8"),
]

for i, (title, desc) in enumerate(feelings):
    y_pos = Inches(2.2) + i * Inches(1.5)
    # Number
    num_shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.1), y_pos + Inches(0.02), Inches(0.35), Inches(0.35))
    num_shape.fill.solid()
    num_shape.fill.fore_color.rgb = ACCENT2
    num_shape.line.fill.background()
    tf = num_shape.text_frame
    p = tf.paragraphs[0]
    p.text = str(i + 1)
    p.font.size = Pt(14)
    p.font.color.rgb = BG_DARK
    p.font.bold = True
    p.font.name = "\ub9d1\uc740 \uace0\ub515"
    p.alignment = PP_ALIGN.CENTER

    add_text(slide, Inches(1.6), y_pos, Inches(4.5), Inches(0.3),
             title, size=16, color=ACCENT2, bold=True)
    add_text(slide, Inches(1.6), y_pos + Inches(0.35), Inches(4.5), Inches(0.9),
             desc, size=12, color=LIGHT)

# Right card - quote
add_shape_bg(slide, Inches(6.7), Inches(1.4), Inches(5.8), Inches(5.2), BG_CARD)
add_text(slide, Inches(7.0), Inches(1.6), Inches(5.2), Inches(0.4),
         "\ud604\uc5c5\uc758 \ubc14\ub78c", size=20, color=YELLOW, bold=True)

# Quote box
add_shape_bg(slide, Inches(7.0), Inches(2.2), Inches(5.2), Inches(3.8), BG_HIGHLIGHT)
add_rect(slide, Inches(7.0), Inches(2.2), Pt(4), Inches(3.8), YELLOW)

quotes = [
    ("\u201c\ud604\uc7a5\uc5d0\uc11c \ubc38\ube0c \ud558\ub098 \ucc3e\ub294 \ub370 10\ubd84 \uac78\ub9ac\ub358 \uac83\uc744,", LIGHT, False, 15),
    ("\uac80\uc0c9 \ud55c \ubc88\uc73c\ub85c \uc904\uc774\uace0 \uc2f6\uc5c8\uc2b5\ub2c8\ub2e4.\u201d", LIGHT, False, 15),
    ("", LIGHT, False, 10),
    ("\u201c\uc815\ube44\uc774\ub825\uc744 \ud55c \ud654\uba74\uc5d0\uc11c \ubcf4\uace0 \uc2f6\ub2e4!\u201d", WHITE, True, 16),
    ("", LIGHT, False, 10),
    ("\u201c\uc88b\uc740 \uac8c \ub9cc\ub4e4\uc5b4\uc84c\uc73c\ub2c8,", LIGHT, False, 15),
    ("\ube68\ub9ac \ud604\uc7a5\uc5d0\uc11c \uc4f0\uace0 \uc2f6\ub2e4!\u201d", WHITE, True, 16),
    ("", LIGHT, False, 10),
    ("\u201cAI\ub85c \uae08\ubc29 \ub9cc\ub4e4\uc5c8\ub294\ub370", LIGHT, False, 15),
    ("\uc65c \ubc14\ub85c \ubabb \uc4f0\uac8c \ud558\ub0d0!\u201d", WHITE, True, 16),
]

add_multiline(slide, Inches(7.3), Inches(2.5), Inches(4.6), Inches(3.2), quotes)


# ============================================================
# SLIDE 6: Governance
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_section_header(slide, "05", "\uac70\ubc84\ub10c\uc2a4 \ud544\uc694\uc131")

add_text(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.5),
         "\ube60\ub974\uac8c \ub9cc\ub4e4 \uc218 \uc788\ub294 \ud798\uc5d0, \uc548\uc2ec\ud558\uace0 \ud65c\uc6a9\ud560 \uc218 \uc788\ub294 \uccb4\uacc4\ub97c \ub354\ud558\uba74 \uc815\ub9d0 \uac15\ub825\ud574\uc9c8 \uac83\uc785\ub2c8\ub2e4",
         size=17, color=YELLOW, bold=True)

gov_items = [
    ("\uc0ac\ub0b4 \uc2dc\uc2a4\ud15c \uc5f0\ub3d9 \uac00\uc774\ub4dc", "EAM, PMS \ub4f1 \ub808\uac70\uc2dc API \uc811\uadfc \uc2dc\n\uc2dc\ud589\ucc29\uc624\ub97c \uc904\uc77c \uc218 \uc788\ub294 \uac00\uc774\ub4dc\ub77c\uc778", ACCENT),
    ("AI \ud504\ub85c\ub355\ud2b8 \uac1c\ubc1c\xb7\ubc30\ud3ec \ud658\uacbd", "\ud604\uc5c5 \ub2f4\ub2f9\uc790\uac00 \uc5b4\ub514\uc5d0\uc11c \uac1c\ubc1c\ud558\uace0\n\uc5b4\ub514\uc5d0 \uc62c\ub9b4\uc9c0 \uace0\ubbfc\ud558\uc9c0 \uc54a\uc544\ub3c4 \ub418\ub294 \ud658\uacbd", ACCENT2),
    ("\ud504\ub85c\ud1a0\ud0c0\uc785 \u2192 \uc6b4\uc601 \uc774\uad00 \uae30\uc900", "\ud604\uc5c5\uc758 \uc18d\ub3c4\ub97c \uc0b4\ub9ac\uba74\uc11c\ub3c4\n\uc548\uc815\uc801\uc73c\ub85c \uc6b4\uc601 \uc774\uad00\ud560 \uc218 \uc788\ub294 \uae30\uc900", GREEN),
    ("\ub370\uc774\ud130 \uac70\ubc84\ub10c\uc2a4", "\ub3c4\uba74, \uc815\ube44\uc774\ub825 \ub4f1 \uc0ac\ub0b4 \ub370\uc774\ud130\ub97c\n\uc5b4\ub514\uc5d0 \uc800\uc7a5\ud558\uace0 \uad00\ub9ac\ud560\uc9c0\uc5d0 \ub300\ud55c \uae30\uc900", ORANGE),
    ("\uc0ac\ub0b4 \uac1c\ubc1c \uc0cc\ub4dc\ubc15\uc2a4", "IT \ud45c\uc900\uacfc \ubcf4\uc548 \ubc14\uc6b4\ub354\ub9ac \uc548\uc5d0\uc11c\n\uc2dc\uc791\ud560 \uc218 \uc788\ub294 \ud50c\ub808\uc774\uadf8\ub77c\uc6b4\ub4dc", RED),
]

for i, (title, desc, color) in enumerate(gov_items):
    x = Inches(0.8) + (i % 3) * Inches(4.0)
    y = Inches(2.1) + (i // 3) * Inches(2.2)
    w = Inches(3.7)
    h = Inches(1.9)
    add_shape_bg(slide, x, y, w, h, BG_CARD)

    num_shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.15), y + Inches(0.2), Inches(0.45), Inches(0.45))
    num_shape.fill.solid()
    num_shape.fill.fore_color.rgb = color
    num_shape.line.fill.background()
    tf = num_shape.text_frame
    p = tf.paragraphs[0]
    p.text = str(i + 1)
    p.font.size = Pt(18)
    p.font.color.rgb = BG_DARK
    p.font.bold = True
    p.font.name = "\ub9d1\uc740 \uace0\ub515"
    p.alignment = PP_ALIGN.CENTER

    add_text(slide, x + Inches(0.7), y + Inches(0.2), w - Inches(0.9), Inches(0.4),
             title, size=18, color=color, bold=True)
    add_text(slide, x + Inches(0.7), y + Inches(0.7), w - Inches(0.9), Inches(1),
             desc, size=13, color=LIGHT)


# ============================================================
# SLIDE 7: 느낀 점 (스토리텔링)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_section_header(slide, "06", "\ub290\ub080 \uc810")

# Opening quote
add_shape_bg(slide, Inches(0.8), Inches(1.4), Inches(11.7), Inches(1.6), BG_HIGHLIGHT)
add_rect(slide, Inches(0.8), Inches(1.4), Pt(4), Inches(1.6), ACCENT)

add_text(slide, Inches(1.3), Inches(1.55), Inches(11), Inches(0.5),
         "\uc800\ub294 \uae30\uacc4 \uc5d4\uc9c0\ub2c8\uc5b4\uc785\ub2c8\ub2e4.", size=22, color=WHITE, bold=True)
add_text(slide, Inches(1.3), Inches(2.05), Inches(11), Inches(0.8),
         "\ucf54\ub529\uc744 \uc804\uacf5\ud558\uc9c0 \uc54a\uc558\uace0, \ud504\ub860\ud2b8\uc5d4\ub4dc \ud504\ub808\uc784\uc6cc\ud06c\ub97c \ubc30\uc6b4 \uc801\ub3c4 \uc5c6\uc2b5\ub2c8\ub2e4.\n\uadf8\ub7f0 \uc81c\uac00 AI\uc640 \ub300\ud654\ud558\uba74\uc11c, 2\uac1c\uc6d4\uac04 48\ubc88\uc758 \ub300\ud654\ub97c \ud1b5\ud574 \ud504\ub85c\ub355\ud2b8\ub97c \ub9cc\ub4e4\uc5c8\uc2b5\ub2c8\ub2e4.",
         size=16, color=LIGHT)

# Three feelings as cards
card_data = [
    ("\uac00\ub2a5\uc131", ACCENT,
     "360\xb0 \ud30c\ub178\ub77c\ub9c8 \ubdf0\uc5b4\ub97c \ub9cc\ub4e4\uc5b4\uc918\ub77c\uace0 \ud588\ub354\ub2c8\nThree.js\ub85c 226\uac1c \uc528\uc744 \uc804\ud658\ud558\ub294 \ucf54\ub4dc\ub97c\n\ub9cc\ub4e4\uc5b4\uc92c\uc2b5\ub2c8\ub2e4.\n\n\uc0c1\uc0c1\ub9cc\uc73c\ub85c \uc2e4\uc81c \ub3d9\uc791\ud558\ub294\n\ud504\ub85c\ub355\ud2b8\ub97c \ub9cc\ub4e4 \uc218 \uc788\uc5c8\uc2b5\ub2c8\ub2e4."),
    ("\uc5f0\uacb0", ACCENT2,
     "GLANCE\uc758 \uc9c4\uc9dc \uac00\uce58\ub294 EAM \uc815\ube44\uc774\ub825\uacfc\n\uc5f0\uacb0\ub420 \ub54c \ub098\uc635\ub2c8\ub2e4.\n\n\uc5f0\ub3d9 \uac00\uc774\ub4dc\uac00 \ubbf8\ub9ac \uc788\uc5c8\ub2e4\uba74\n\uc2dc\ud589\ucc29\uc624\ub97c \ud06c\uac8c \uc904\uc77c \uc218 \uc788\uc5c8\uc744 \uac83\uc785\ub2c8\ub2e4."),
    ("\ud65c\uc6a9", ORANGE,
     "\uac1c\ubc1c\uc740 \uc26c\uc6cc\uc84c\ub294\ub370\n\uadf8 \ub2e4\uc74c \ub2e8\uacc4\uc758 \uae38\uc774 \ubcf4\uc774\uc9c0 \uc54a\uc558\uc2b5\ub2c8\ub2e4.\n\n\uc5b4\ub514\uc5d0 \ubc30\ud3ec\ud574\uc57c \ud558\ub294\uc9c0,\n\uc5b4\ub5a4 \uc694\uac74\uc744 \uac16\ucdb0\uc57c \ud558\ub294\uc9c0\n\uae30\uc900\uc774 \uc788\uc5c8\uc73c\uba74 \uc88b\uaca0\ub2e4\uace0 \ub290\uaf08\uc2b5\ub2c8\ub2e4."),
]

for i, (label, color, desc) in enumerate(card_data):
    x = Inches(0.8) + i * Inches(4.0)
    y = Inches(3.4)
    w = Inches(3.7)
    h = Inches(3.5)
    add_shape_bg(slide, x, y, w, h, BG_CARD)
    # Color bar top
    add_rect(slide, x, y, w, Pt(4), color)
    add_text(slide, x + Inches(0.25), y + Inches(0.25), w - Inches(0.5), Inches(0.4),
             label, size=20, color=color, bold=True)
    add_text(slide, x + Inches(0.25), y + Inches(0.75), w - Inches(0.5), Inches(2.5),
             desc, size=13, color=LIGHT)


# ============================================================
# SLIDE 8: 제언
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_section_header(slide, "07", "\uc81c\uc5b8")

# Main message
add_shape_bg(slide, Inches(0.8), Inches(1.4), Inches(11.7), Inches(1.2), BG_HIGHLIGHT)
add_rect(slide, Inches(0.8), Inches(1.4), Pt(4), Inches(1.2), YELLOW)
add_text(slide, Inches(1.3), Inches(1.55), Inches(11), Inches(0.4),
         "Pio\uc758 Blind \uad00\ub9ac \uc2dc\uc2a4\ud15c\uc774\ub4e0, GLANCE\ub4e0, \ubb38\uc81c\uc758 \ubcf8\uc9c8\uc740 \uac19\uc2b5\ub2c8\ub2e4.",
         size=18, color=WHITE, bold=True)
add_text(slide, Inches(1.3), Inches(1.95), Inches(11), Inches(0.4),
         "AI\ub85c \ube60\ub974\uac8c \ub9cc\ub4e4 \uc218 \uc788\uac8c \ub410\ub2e4. \uadf8\ub7f0\ub370 \uadf8 \ub2e4\uc74c \ub2e8\uacc4\uac00 \ub9c9\ud78c\ub2e4.",
         size=16, color=YELLOW)

# Table
add_text(slide, Inches(0.8), Inches(2.9), Inches(11.7), Inches(0.4),
         "\ubbf8\ub9ac \uc900\ube44\ub418\uc5b4 \uc788\uc73c\uba74 \uc88b\uc558\uc744 \uac83\ub4e4", size=18, color=WHITE, bold=True)

table_items = [
    ("\uc0ac\ub0b4 \uc2dc\uc2a4\ud15c \uc5f0\ub3d9 \uac00\uc774\ub4dc", "EAM, PMS \ub4f1 \ub808\uac70\uc2dc API \uc811\uadfc \uc2dc \uc2dc\ud589\ucc29\uc624 \uac10\uc18c"),
    ("AI \ud504\ub85c\ub355\ud2b8 \uac1c\ubc1c\xb7\ubc30\ud3ec \ud658\uacbd", "\uc5b4\ub514\uc5d0\uc11c \uac1c\ubc1c\ud558\uace0 \uc5b4\ub514\uc5d0 \uc62c\ub9b4\uc9c0 \uace0\ubbfc\ud558\uc9c0 \uc54a\uc544\ub3c4 \ub418\ub294 \ud658\uacbd"),
    ("\ud504\ub85c\ud1a0\ud0c0\uc785 \u2192 \uc6b4\uc601 \uc774\uad00 \uae30\uc900", "\ud604\uc5c5\uc758 \uc18d\ub3c4\ub97c \uc0b4\ub9ac\uba74\uc11c \uc548\uc815\uc801\uc73c\ub85c \uc6b4\uc601 \uc804\ud658"),
]

for i, (item, effect) in enumerate(table_items):
    y = Inches(3.4) + i * Inches(0.7)
    # Left
    add_shape_bg(slide, Inches(0.8), y, Inches(4.5), Inches(0.55), BG_CARD)
    add_text(slide, Inches(1.0), y + Inches(0.1), Inches(4.1), Inches(0.4),
             item, size=15, color=ACCENT, bold=True)
    # Arrow
    add_text(slide, Inches(5.4), y + Inches(0.05), Inches(0.5), Inches(0.4),
             "\u2192", size=18, color=GRAY, align=PP_ALIGN.CENTER)
    # Right
    add_shape_bg(slide, Inches(6.0), y, Inches(6.5), Inches(0.55), BG_CARD)
    add_text(slide, Inches(6.2), y + Inches(0.1), Inches(6.1), Inches(0.4),
             effect, size=14, color=LIGHT)

# Bottom highlight
add_shape_bg(slide, Inches(0.8), Inches(5.6), Inches(11.7), Inches(1.2), BG_HIGHLIGHT)
add_rect(slide, Inches(0.8), Inches(5.6), Pt(4), Inches(1.2), GREEN)
add_text(slide, Inches(1.3), Inches(5.75), Inches(11), Inches(0.4),
         "\uc774\ub7f0 \uccb4\uacc4\uac00 \uac16\ucdb0\uc9c4\ub2e4\uba74,", size=18, color=WHITE, bold=True)
add_text(slide, Inches(1.3), Inches(6.15), Inches(11), Inches(0.5),
         "\ud604\uc5c5\uc758 \uc544\uc774\ub514\uc5b4\uac00 \ud504\ub85c\ud1a0\ud0c0\uc785\uc5d0\uc11c \uba48\ucd94\uc9c0 \uc54a\uace0, \uc2e4\uc81c \uc5c5\ubb34 \ud604\uc7a5\uae4c\uc9c0 \ube60\ub974\uac8c \ub3c4\ub2ec\ud560 \uc218 \uc788\uc744 \uac83\uc785\ub2c8\ub2e4.",
         size=16, color=GREEN)


# ============================================================
# SLIDE 9: Closing
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_rect(slide, Inches(0), Inches(3.6), Inches(13.333), Pt(3), ACCENT)

add_text(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(0.7),
         "AI Code Generation\uc740", size=24, color=GRAY, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1.5), Inches(2.1), Inches(10), Inches(0.8),
         "\u201c\ub9cc\ub4dc\ub294 \uc0ac\ub78c\u201d\uc758 \ubc94\uc704\ub97c \ub113\ud600\uc92c\uc2b5\ub2c8\ub2e4.", size=30, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_text(slide, Inches(1.5), Inches(3.0), Inches(10), Inches(0.5),
         "\ub354 \uc774\uc0c1 \uac1c\ubc1c\uc790\ub9cc \uc18c\ud504\ud2b8\uc6e8\uc5b4\ub97c \ub9cc\ub4dc\ub294 \uc2dc\ub300\uac00 \uc544\ub2d9\ub2c8\ub2e4.", size=18, color=GRAY, align=PP_ALIGN.CENTER)

add_text(slide, Inches(1.5), Inches(4.1), Inches(10), Inches(0.7),
         "\ube60\ub974\uac8c \ub9cc\ub4e4 \uc218 \uc788\ub294 \uc774 \ud798\uc744,\n\uc870\uc9c1 \uc548\uc5d0\uc11c \uc548\uc2ec\ud558\uace0 \ud65c\uc6a9\ud560 \uc218 \uc788\ub294 \uccb4\uacc4\ub97c \ud568\uaed8 \ub9cc\ub4e4\uc5b4\uac11\uc2dc\ub2e4.",
         size=20, color=ACCENT, align=PP_ALIGN.CENTER)

add_text(slide, Inches(1.5), Inches(5.5), Inches(10), Inches(0.5),
         "\uac10\uc0ac\ud569\ub2c8\ub2e4", size=28, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_text(slide, Inches(1.5), Inches(6.2), Inches(10), Inches(0.4),
         "GLANCE \xb7 AI Code Generation \xb7 52G \ud611\uc758\uccb4", size=14, color=GRAY, align=PP_ALIGN.CENTER)


# Save
output_path = r"c:\Users\USER\Downloads\pid-viewer\GLANCE_AI_Governance_v2.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
